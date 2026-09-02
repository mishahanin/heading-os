"""An extractor that could not read most of the corpus it was pointed at.

A datastore document reaches the memory index, and therefore `/recall`, only
through its `-extract.md` companion. `scripts/datastore-extract.py` knew two
formats, `.xlsx` and `.pptx`, so a `.pdf` or a `.docx` had no companion, could
not have one, and was unreachable by every retrieval path in the workspace.
Measured in the operator's datastore on 2026-09-02: 165 of 983 files were PDF
or DOCX and could never be reached at all.

This file covers the two new extractors and the three judgement calls they
carry, plus a control over the two old ones, because the cheapest way to add a
format is to change the shared scan loop underneath the formats that worked.

Fixtures are built here at test time rather than committed: a PDF is written as
bytes by `_build_pdf` below (no pure-python PDF writer is installed, see its
docstring), and the DOCX, XLSX and PPTX fixtures come from the same libraries
the extractors read them back with.

Tests: this file. Sibling coverage of the same script:
tests/test_a_rollback_that_deleted_what_it_never_backed_up.py
"""
from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))


def _load(stem: str, name: str):
    spec = importlib.util.spec_from_file_location(name, ROOT / "scripts" / f"{stem}.py")
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    spec.loader.exec_module(mod)
    return mod


extract = _load("datastore-extract", "datastore_extract_pdf_docx")


# ==========================================================================
# Fixtures: everything is built in tmp_path, nothing is committed
# ==========================================================================

@pytest.fixture(autouse=True)
def _never_the_real_overlay(tmp_path, monkeypatch):
    """Repoint the data root at tmp_path for every test in this module.

    Belt and braces. Each test below also hands `scan_report` an explicit
    target directory, so none of them resolves the data root at all; this
    fixture is what makes a future test that FORGETS to do so land in tmp_path
    instead of the operator's real overlay. `datastore_dir()` reads the
    environment on every call precisely so that this works, and
    `test_the_scan_target_follows_the_environment` below pins that.
    """
    data_root = tmp_path / "data-root"
    (data_root / "datastore").mkdir(parents=True)
    monkeypatch.setenv("HEADING_OS_DATA", str(data_root))


def _build_pdf(page_streams: list, with_font: bool = True) -> bytes:
    """A structurally valid PDF, one page per content stream.

    Hand-assembled bytes, with the cross-reference offsets computed rather than
    faked, because no pure-python PDF writer is importable: measured
    2026-09-02, `importlib.util.find_spec` returns None in this project's
    `.venv` for every one of reportlab, fpdf, pypdf, PyPDF2, pikepdf,
    weasyprint and borb. If any of those is ever added, this helper can go.

    `with_font=False` omits the font resource, which is how the no-text-layer
    fixture is built: a page that draws a filled rectangle and nothing else.
    """
    n_pages = len(page_streams)
    font_num = 3 + 2 * n_pages
    kids = " ".join(f"{3 + 2 * i} 0 R" for i in range(n_pages))

    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        f"<< /Type /Pages /Kids [{kids}] /Count {n_pages} >>".encode(),
    ]
    for i, stream in enumerate(page_streams):
        resources = f"/Font << /F1 {font_num} 0 R >>" if with_font else ""
        objects.append(
            f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            f"/Resources << {resources} >> /Contents {4 + 2 * i} 0 R >>".encode())
        data = stream.encode("latin-1")
        objects.append(b"<< /Length %d >>\nstream\n" % len(data) + data + b"\nendstream")
    if with_font:
        objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for number, body in enumerate(objects, 1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % number + body + b"\nendobj\n"

    xref_at = len(out)
    out += b"xref\n0 %d\n" % (len(objects) + 1)
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += b"%010d 00000 n \n" % offset
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n" % (
        len(objects) + 1, xref_at)
    return bytes(out)


def _text_page(lines: list) -> str:
    """A content stream that writes the given lines with the base-14 font."""
    parts = ["BT", "/F1 12 Tf", "72 720 Td", "14 TL"]
    for line in lines:
        escaped = line.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
        parts.append(f"({escaped}) Tj T*")
    parts.append("ET")
    return "\n".join(parts)


def _write_text_pdf(path: Path, pages: list) -> Path:
    path.write_bytes(_build_pdf([_text_page(lines) for lines in pages]))
    return path


def _write_scanned_pdf(path: Path) -> Path:
    """A valid one-page PDF with no text operators at all."""
    path.write_bytes(_build_pdf(["0 0 0 rg 100 100 200 200 re f"], with_font=False))
    return path


def _write_docx(path: Path, paragraphs: list, table: list | None = None) -> Path:
    """A DOCX whose table, when present, sits BETWEEN the paragraphs."""
    import docx

    document = docx.Document()
    if paragraphs:
        document.add_paragraph(paragraphs[0])
    if table:
        rendered = document.add_table(rows=len(table), cols=len(table[0]))
        for r, row in enumerate(table):
            for c, value in enumerate(row):
                rendered.cell(r, c).text = value
    for text in paragraphs[1:]:
        document.add_paragraph(text)
    document.save(str(path))
    return path


def _write_xlsx(path: Path) -> Path:
    from openpyxl import Workbook

    book = Workbook()
    sheet = book.active
    sheet.title = "Coverage"
    sheet.append(["Region", "Sensors"])
    sheet.append(["North basin", 12])
    book.save(str(path))
    return path


def _write_pptx(path: Path) -> Path:
    from pptx import Presentation

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Heading for the quarter"
    deck.save(str(path))
    return path


# ==========================================================================
# 1 - the two suffixes the scan could not see
# ==========================================================================

@pytest.mark.parametrize("suffix", [".pdf", ".docx"])
def test_the_missing_suffixes_are_now_extractable(suffix):
    assert suffix in extract.EXTRACTABLE_SUFFIXES, (
        f"{suffix} files have no companion and no way to get one, so nothing "
        f"in the memory index can reach their contents")


def test_the_original_suffixes_are_still_extractable():
    """The control on the tuple itself: widening it must not drop anything."""
    assert ".xlsx" in extract.EXTRACTABLE_SUFFIXES
    assert ".pptx" in extract.EXTRACTABLE_SUFFIXES


def test_a_pdf_and_a_docx_are_picked_up_by_a_scan(tmp_path):
    _write_text_pdf(tmp_path / "field-notes.pdf", [["A line of readable text"]])
    _write_docx(tmp_path / "rate-card.docx", ["A readable paragraph"])

    extract.scan_report(target_dir=tmp_path)

    assert (tmp_path / "field-notes-extract.md").is_file()
    assert (tmp_path / "rate-card-extract.md").is_file()


# ==========================================================================
# 2 - the DOCX whose content is all in tables
# ==========================================================================

def test_a_docx_table_reaches_the_extract(tmp_path):
    """`Document.paragraphs` and `Document.tables` are separate collections.

    Reading only the first is the plausible implementation, and it returns
    nothing at all for the documents that are mostly grids.
    """
    path = _write_docx(tmp_path / "spec.docx", ["Intro paragraph"],
                       table=[["Module", "Seats"], ["Sensor", "4150"]])
    text = extract.extract_docx(path)

    assert "Intro paragraph" in text
    for value in ("Module", "Seats", "Sensor", "4150"):
        assert value in text, f"the table cell {value!r} was never extracted"


def test_a_docx_that_is_only_a_table_is_not_empty(tmp_path):
    path = _write_docx(tmp_path / "matrix.docx", [],
                       table=[["Field", "Value"], ["Serial", "AX-9931"]])
    text = extract.extract_docx(path)

    assert "AX-9931" in text, "a document whose content is all table came out empty"
    assert "No extractable text layer" not in text, (
        "a table-only document was reported as having no text layer")


def test_a_docx_table_is_rendered_as_a_markdown_table(tmp_path):
    path = _write_docx(tmp_path / "grid.docx", ["Intro"],
                       table=[["Module", "Seats"], ["Sensor", "4150"]])
    text = extract.extract_docx(path)

    assert "| Module | Seats |" in text
    assert "| --- | --- |" in text


def test_a_pipe_in_a_table_cell_does_not_add_a_column(tmp_path):
    """The same escaping `_cell` already gives XLSX rows."""
    path = _write_docx(tmp_path / "piped.docx", ["Intro"],
                       table=[["Module", "Rate | per seat"], ["Sensor", "4150"]])
    text = extract.extract_docx(path)

    assert r"Rate \| per seat" in text, (
        "a raw pipe in a cell added a phantom column to the companion")


def test_a_docx_keeps_the_document_order(tmp_path):
    """Paragraphs then tables would put the closing paragraph before the grid."""
    path = _write_docx(tmp_path / "ordered.docx", ["Opening line", "Closing line"],
                       table=[["Module", "Seats"], ["Sensor", "4150"]])
    text = extract.extract_docx(path)

    assert text.index("Opening line") < text.index("| Module | Seats |")
    assert text.index("| Module | Seats |") < text.index("Closing line")


# ==========================================================================
# 3 - the scan with no text layer, which must not read as an empty document
# ==========================================================================

def test_a_scanned_pdf_says_it_has_no_text_layer(tmp_path):
    path = _write_scanned_pdf(tmp_path / "signed-page.pdf")
    text = extract.extract_pdf(path)

    assert "No extractable text layer" in text
    assert "NOT the record of an empty document" in text, (
        "the companion for a scan does not say what it is, so the next reader "
        "takes it for a document with nothing in it")
    assert "signed-page.pdf" in text, "the companion never names the binary to open"


def test_a_scanned_pdf_companion_is_more_than_its_header(tmp_path):
    """The failure this guards is a companion that is only the header block.

    A header-only file is the exact shape of "this document is empty", so the
    body is compared against the header the other extractors also write.
    """
    path = _write_scanned_pdf(tmp_path / "signed-page.pdf")
    text = extract.extract_pdf(path)
    header = "\n".join(extract._extract_header(path, "PDF"))

    assert text.startswith(header), "the shared header convention was dropped"
    body = text[len(header):].strip()
    assert len(body) > 100, (
        f"the companion is {len(body)} characters past its header, which reads "
        f"as an empty document rather than an unreadable one")


def test_a_scanned_pdf_is_written_rather_than_skipped(tmp_path, capsys):
    """The documented choice: write a companion, do not skip and tally."""
    _write_scanned_pdf(tmp_path / "signed-page.pdf")

    result = extract.scan_report(target_dir=tmp_path)

    companion = tmp_path / "signed-page-extract.md"
    assert companion.is_file(), "the scan left no trace of the file on disk"
    assert "No extractable text layer" in companion.read_text(encoding="utf-8")
    assert result.failures == [], "a scan is unreadable, not a failure to report"


def test_a_docx_with_no_content_says_so_too(tmp_path):
    """The DOCX half of the same decision.

    Measured 2026-09-02: with only the PDF case covered, deleting the DOCX
    no-text branch left all 34 tests in this file green, and an empty document
    got a companion consisting of a header, `Paragraphs: 0` and `Tables: 0`.
    """
    import docx

    path = tmp_path / "blank.docx"
    docx.Document().save(str(path))

    text = extract.extract_docx(path)

    assert "No extractable text layer" in text
    assert "NOT the record of an empty document" in text


def test_a_pdf_with_text_is_not_marked_as_having_none(tmp_path):
    """Negative case, so the assertion above is not satisfied by everything."""
    path = _write_text_pdf(tmp_path / "brief.pdf", [["Readable body text here"]])
    text = extract.extract_pdf(path)

    assert "No extractable text layer" not in text
    assert "Readable body text here" in text


# ==========================================================================
# 4 - the corrupt file that must be named, counted, and survived
# ==========================================================================

def test_a_corrupt_pdf_is_counted_as_a_failure(tmp_path):
    (tmp_path / "truncated.pdf").write_bytes(b"%PDF-1.4\nthis is not a pdf\n")

    result = extract.scan_report(target_dir=tmp_path)

    assert len(result.failures) == 1, "an unreadable PDF passed silently"
    path, why = result.failures[0]
    assert path.name == "truncated.pdf"
    assert "PDFSyntaxError" in why, f"the failure was not named, only counted: {why}"


def test_a_corrupt_pdf_does_not_stop_the_batch(tmp_path):
    """`truncated.pdf` sorts first, so the good file comes after the bad one."""
    (tmp_path / "truncated.pdf").write_bytes(b"%PDF-1.4\nthis is not a pdf\n")
    _write_docx(tmp_path / "zz-later.docx", ["This paragraph must still arrive"])

    result = extract.scan_report(target_dir=tmp_path)

    assert (tmp_path / "zz-later-extract.md").is_file(), (
        "one unreadable file ended the loop and every later file went unread")
    assert len(result.extracted) == 1
    assert not (tmp_path / "truncated-extract.md").exists(), (
        "a companion was written for a file that could not be parsed")


def test_a_corrupt_pdf_is_named_in_the_report(tmp_path, capsys):
    (tmp_path / "truncated.pdf").write_bytes(b"%PDF-1.4\nthis is not a pdf\n")

    extract.scan_report(target_dir=tmp_path)

    out = capsys.readouterr().out
    assert "1 file(s) could not be extracted" in out
    assert "truncated.pdf" in out


def test_the_exit_code_reflects_a_failure(tmp_path, monkeypatch):
    (tmp_path / "truncated.pdf").write_bytes(b"%PDF-1.4\nthis is not a pdf\n")
    _write_docx(tmp_path / "zz-later.docx", ["Extracted fine"])
    monkeypatch.setattr(sys, "argv", ["datastore-extract.py", str(tmp_path)])

    assert extract.main() == 1, (
        "the run reported success in its exit status while one of its input "
        "files was never read")


def test_the_exit_code_is_zero_with_no_failures(tmp_path, monkeypatch):
    _write_docx(tmp_path / "clean.docx", ["Extracted fine"])
    monkeypatch.setattr(sys, "argv", ["datastore-extract.py", str(tmp_path)])

    assert extract.main() == 0


def test_an_encrypted_pdf_is_a_failure_not_a_silent_pass(tmp_path):
    """pdfminer raises out of `extract_pdf`; nothing in it may swallow that."""
    pdf = _build_pdf([_text_page(["secret"])]).replace(
        b"/Root 1 0 R", b"/Encrypt 99 0 R /Root 1 0 R")
    (tmp_path / "locked.pdf").write_bytes(pdf)

    result = extract.scan_report(target_dir=tmp_path)

    assert result.extracted == []
    assert len(result.failures) == 1, "an unreadable PDF was treated as extracted"
    # Named, so this cannot start passing because the fixture became merely
    # malformed. Measured 2026-09-02: pdfminer raises PDFEncryptionError
    # ("Unknown filter") on this file, not PDFSyntaxError.
    assert "PDFEncryptionError" in result.failures[0][1], (
        f"the encrypted file failed for some other reason: "
        f"{result.failures[0][1]}")


# ==========================================================================
# 5 - the cap, and the mark that says the companion is partial
# ==========================================================================

def test_the_cap_is_the_documented_number():
    """The docstring states 100,000 characters and the reason for it.

    Pinned so the prose and the constant cannot drift apart silently; a
    deliberate change to the cap edits both.
    """
    assert extract.MAX_EXTRACT_CHARS == 100_000


def test_a_long_pdf_is_truncated_at_a_page_boundary(tmp_path, monkeypatch):
    monkeypatch.setattr(extract, "MAX_EXTRACT_CHARS", 40)
    path = _write_text_pdf(tmp_path / "manual.pdf", [
        ["First page body text"],
        ["Second page body text"],
        ["Third page body text"],
    ])

    text = extract.extract_pdf(path)

    assert "First page body text" in text
    assert "Third page body text" not in text, "the cap did not stop the extract"
    assert "## Page 3" not in text, "the cut landed inside a page, not between two"


def test_a_truncated_pdf_says_so_at_both_ends(tmp_path, monkeypatch):
    """A reader who sees only the head of the file must still be told.

    The memory index chunks a companion, so a marker written only at the point
    of the cut reaches whoever reads the tail and nobody else.

    The marker LINES are located rather than a leading and trailing slice
    searched: on a fixture this small the two slices overlap, so `"TRUNCATED"
    in head` and `"TRUNCATED" in tail` were both satisfied by a single marker
    at either end. Measured 2026-09-02, deleting either of the two writes left
    all 34 tests in this file green.
    """
    monkeypatch.setattr(extract, "MAX_EXTRACT_CHARS", 40)
    path = _write_text_pdf(tmp_path / "manual.pdf", [
        ["First page body text"], ["Second page body text"], ["Third page body text"],
    ])

    lines = extract.extract_pdf(path).splitlines()
    marks = [i for i, line in enumerate(lines) if "TRUNCATED" in line]
    first_page = next(i for i, line in enumerate(lines) if line.startswith("## Page"))

    assert len(marks) == 2, (
        f"the partial companion carries {len(marks)} truncation marker(s); a "
        f"reader who gets only one end of it is told nothing")
    assert marks[0] < first_page, "no marker above the content"
    assert marks[1] > first_page, "no marker below the content"
    assert "PARTIAL" in lines[marks[0]]
    assert "40-character cap" in lines[marks[0]], (
        "the marker does not say what stopped the extract")


def test_a_short_pdf_carries_no_truncation_mark(tmp_path):
    """Negative case for the two assertions above."""
    path = _write_text_pdf(tmp_path / "short.pdf", [["One short page"]])

    assert "TRUNCATED" not in extract.extract_pdf(path)


def test_a_long_docx_is_truncated_at_a_block_boundary(tmp_path, monkeypatch):
    monkeypatch.setattr(extract, "MAX_EXTRACT_CHARS", 30)
    path = _write_docx(tmp_path / "long.docx", [
        "A paragraph of about forty characters.",
        "A second paragraph that must be cut off.",
    ])

    text = extract.extract_docx(path)

    assert "A paragraph of about forty characters." in text, (
        "the cut landed inside the first block instead of after it")
    assert "must be cut off" not in text
    assert "TRUNCATED" in text and "PARTIAL" in text


def test_a_short_docx_carries_no_truncation_mark(tmp_path):
    path = _write_docx(tmp_path / "short.docx", ["One short paragraph"])

    assert "TRUNCATED" not in extract.extract_docx(path)


# ==========================================================================
# 6 - the header convention the new extractors had to join
# ==========================================================================

@pytest.mark.parametrize("name,writer,kind", [
    ("brief.pdf", lambda p: _write_text_pdf(p, [["Body text"]]), "PDF"),
    ("brief.docx", lambda p: _write_docx(p, ["Body text"]), "DOCX"),
])
def test_the_new_companions_carry_the_shared_header(tmp_path, name, writer, kind):
    from datetime import datetime

    from scripts.utils.workspace import get_default_tz

    path = writer(tmp_path / name)
    text = getattr(extract, f"extract_{kind.lower()}")(path)
    today = datetime.now(get_default_tz()).strftime("%Y-%m-%d")

    assert text.startswith(f"# Extract: {name}\n")
    assert f"> Auto-extracted from `{name}` on {today}" in text
    assert f"The original {kind} is the source of truth." in text


# ==========================================================================
# 7 - the control: the two formats that already worked
# ==========================================================================

def test_an_xlsx_extract_is_unchanged(tmp_path):
    path = _write_xlsx(tmp_path / "coverage.xlsx")
    text = extract.extract_xlsx(path)

    assert text.startswith("# Extract: coverage.xlsx\n")
    assert "The original XLSX is the source of truth." in text
    assert "## Sheet: Coverage" in text
    assert "| Region | Sensors |" in text
    assert "| North basin | 12 |" in text


def test_a_pptx_extract_is_unchanged(tmp_path):
    path = _write_pptx(tmp_path / "deck.pptx")
    text = extract.extract_pptx(path)

    assert text.startswith("# Extract: deck.pptx\n")
    assert "The original PPTX is the source of truth." in text
    assert "Total slides: 1" in text
    assert "## Slide 1" in text
    assert "- Heading for the quarter" in text


def test_all_four_formats_extract_in_one_scan(tmp_path):
    _write_xlsx(tmp_path / "coverage.xlsx")
    _write_pptx(tmp_path / "deck.pptx")
    _write_text_pdf(tmp_path / "brief.pdf", [["Body text"]])
    _write_docx(tmp_path / "notes.docx", ["Body text"])

    result = extract.scan_report(target_dir=tmp_path)

    assert {orig.name for orig, _ in result.extracted} == {
        "coverage.xlsx", "deck.pptx", "brief.pdf", "notes.docx"}
    assert result.failures == []


def test_scan_and_extract_still_returns_a_bare_list(tmp_path):
    """Three other test modules and `update_index` consume exactly this.

    `scan_report` is the widened sibling; `scan_and_extract` keeps the old
    contract so adding the failure channel could not break them.
    """
    _write_docx(tmp_path / "notes.docx", ["Body text"])

    extracted = extract.scan_and_extract(tmp_path)

    assert isinstance(extracted, list)
    assert [orig.name for orig, _ in extracted] == ["notes.docx"]


def test_an_unrelated_extension_is_still_ignored(tmp_path, capsys):
    (tmp_path / "notes.txt").write_text("x", encoding="utf-8")

    extract.scan_report(target_dir=tmp_path)

    assert "No XLSX or PPTX files" in capsys.readouterr().out, (
        "the message the sibling test module asserts on was reworded")


def test_a_pdf_and_an_xlsx_sharing_a_stem_get_distinct_companions(tmp_path):
    """The collision rule was written for a .pptx/.xlsx pair; there are now six."""
    _write_xlsx(tmp_path / "pitch.xlsx")
    _write_text_pdf(tmp_path / "pitch.pdf", [["Body text"]])

    extract.scan_report(target_dir=tmp_path)

    assert (tmp_path / "pitch-xlsx-extract.md").is_file()
    assert (tmp_path / "pitch-pdf-extract.md").is_file()
    assert not (tmp_path / "pitch-extract.md").exists(), (
        "both sources wrote to one companion and the second overwrote the first")


# ==========================================================================
# 8 - the frozen constant that once wrote the operator's real overlay
# ==========================================================================

def test_the_scan_target_follows_the_environment(monkeypatch, tmp_path):
    """`datastore_dir()` resolves per call, never once at import.

    The docstring on `datastore_dir` records why: as a module-level constant it
    answered during its own import, so a test that repointed the data root
    afterwards still read and wrote the operator's real datastore. Adding two
    formats multiplies what such a run would touch, so the property is pinned
    here as well.
    """
    first = tmp_path / "one"
    second = tmp_path / "two"
    monkeypatch.setattr(extract, "get_datastore_dir", lambda: first)
    assert extract.datastore_dir() == first

    monkeypatch.setattr(extract, "get_datastore_dir", lambda: second)
    assert extract.datastore_dir() == second, (
        "the datastore path was frozen at import and no longer follows the "
        "environment the caller set")
