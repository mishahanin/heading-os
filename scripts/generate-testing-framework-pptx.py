#!/usr/bin/env python3
"""Generate MIB Testing Framework presentation in 31C corporate style."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from scripts.utils.venv import ensure_venv  # noqa: E402

ensure_venv()
from scripts.utils.workspace import get_outputs_dir

# 31C Brand Colors
DARK_BG = RGBColor(0x0A, 0x0E, 0x1A)
CARD_BG = RGBColor(0x12, 0x17, 0x2A)
BRAND_BLUE = RGBColor(0x42, 0x3B, 0xFF)
BRAND_ORANGE = RGBColor(0xFF, 0x92, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
MED_GRAY = RGBColor(0x6B, 0x72, 0x85)
ACCENT_TEAL = RGBColor(0x00, 0xD4, 0xAA)
CARD_BORDER = RGBColor(0x1E, 0x24, 0x3B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUTPUT = str(get_outputs_dir() / "deliverables" / "presentations"
             / "MIB - Testing Framework (31C Style).pptx")


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill=None, border=None, bw=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or CARD_BG
    if border:
        shape.line.color.rgb = border
        shape.line.width = bw or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_bar(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def txt(slide, left, top, w, h, text, size=14, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font="Segoe UI", spacing=1.2):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    if spacing != 1.0:
        p.line_spacing = Pt(size * spacing)
    return box


def top_bar(slide):
    add_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BRAND_BLUE)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # =========================================================
    # SLIDE 1: TITLE
    # =========================================================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s1, DARK_BG)
    top_bar(s1)

    txt(s1, Inches(0.8), Inches(1.5), Inches(4), Inches(0.5),
        "31 CONCEPT", size=14, color=BRAND_BLUE, bold=True)

    txt(s1, Inches(0.8), Inches(2.2), Inches(8), Inches(1.4),
        "Testing Framework\n& Strategic Advantages",
        size=44, color=WHITE, bold=True, spacing=1.1)

    add_bar(s1, Inches(0.8), Inches(3.9), Inches(2.5), Pt(3), BRAND_ORANGE)

    txt(s1, Inches(0.8), Inches(4.15), Inches(8), Inches(0.8),
        "Comprehensive evaluation methodology for\nODUN.ONE platform validation",
        size=16, color=LIGHT_GRAY)

    txt(s1, Inches(0.8), Inches(6.5), Inches(4), Inches(0.4),
        "PROPRIETARY & CONFIDENTIAL", size=10, color=MED_GRAY)

    # Decorative geometry (right)
    add_rect(s1, Inches(9.5), Inches(1.0), Inches(3.5), Inches(3.5),
             fill=RGBColor(0x14, 0x18, 0x30), border=BRAND_BLUE, bw=Pt(1.5))
    add_rect(s1, Inches(10.0), Inches(1.5), Inches(2.8), Inches(2.8),
             fill=RGBColor(0x18, 0x1D, 0x38), border=RGBColor(0x2A, 0x25, 0xCC))

    # =========================================================
    # SLIDE 2: MAIN PURPOSE (5 numbered cards)
    # =========================================================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s2, DARK_BG)
    top_bar(s2)

    txt(s2, Inches(0.8), Inches(0.4), Inches(3), Inches(0.4),
        "01 / MAIN PURPOSE", size=12, color=BRAND_ORANGE, bold=True)
    txt(s2, Inches(0.8), Inches(0.85), Inches(10), Inches(0.6),
        "What We Set Out to Prove", size=32, color=WHITE, bold=True)

    purposes = [
        ("Benchmark Platform Capability",
         "Evaluate functional and non-functional performance against key competitors."),
        ("Concept Validation",
         "High-level capability testing \u2014 not production-grade validation."),
        ("Evaluate Team Performance",
         "Measure professionalism, responsiveness, and adaptability under test conditions."),
        ("Use Case Feasibility",
         "Determine technical feasibility for executing real-world use cases."),
        ("End-to-End Assessment",
         "Assess the platform end-to-end \u2014 not limited to specific use cases."),
    ]

    cw, ch = Inches(3.7), Inches(2.2)
    gap = Inches(0.3)
    sy = Inches(1.75)

    # Row 1: 3 cards
    for i in range(3):
        x = Inches(0.6) + i * (cw + gap)
        add_rect(s2, x, sy, cw, ch, fill=CARD_BG, border=CARD_BORDER)
        add_rect(s2, x + Inches(0.25), sy + Inches(0.25), Inches(0.45), Inches(0.45),
                 fill=BRAND_BLUE)
        txt(s2, x + Inches(0.25), sy + Inches(0.28), Inches(0.45), Inches(0.4),
            str(i + 1), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s2, x + Inches(0.85), sy + Inches(0.3), cw - Inches(1.1), Inches(0.4),
            purposes[i][0], size=14, color=WHITE, bold=True)
        txt(s2, x + Inches(0.25), sy + Inches(1.0), cw - Inches(0.5), Inches(1.0),
            purposes[i][1], size=11.5, color=LIGHT_GRAY)

    # Row 2: 2 cards centered
    rx = Inches(0.6) + (cw + gap) * 0.5
    for i in range(2):
        idx = i + 3
        x = rx + i * (cw + gap)
        y = sy + ch + gap
        add_rect(s2, x, y, cw, ch, fill=CARD_BG, border=CARD_BORDER)
        add_rect(s2, x + Inches(0.25), y + Inches(0.25), Inches(0.45), Inches(0.45),
                 fill=BRAND_BLUE)
        txt(s2, x + Inches(0.25), y + Inches(0.28), Inches(0.45), Inches(0.4),
            str(idx + 1), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s2, x + Inches(0.85), y + Inches(0.3), cw - Inches(1.1), Inches(0.4),
            purposes[idx][0], size=14, color=WHITE, bold=True)
        txt(s2, x + Inches(0.25), y + Inches(1.0), cw - Inches(0.5), Inches(1.0),
            purposes[idx][1], size=11.5, color=LIGHT_GRAY)

    # =========================================================
    # SLIDE 3: STRATEGIC ADVANTAGES (two-column cards)
    # =========================================================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s3, DARK_BG)
    top_bar(s3)

    txt(s3, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
        "02 / STRATEGIC ADVANTAGES", size=12, color=BRAND_ORANGE, bold=True)
    txt(s3, Inches(0.8), Inches(0.85), Inches(10), Inches(0.6),
        "Why 31 Concept Wins", size=32, color=WHITE, bold=True)

    def card_col(slide, x, y, w, h, heading, hcolor, items):
        add_rect(slide, x, y, w, h, fill=CARD_BG, border=CARD_BORDER)
        # Top accent bar
        add_bar(slide, x, y, w, Pt(4), hcolor)
        # Heading
        txt(slide, x + Inches(0.35), y + Inches(0.3), w - Inches(0.7), Inches(0.5),
            heading, size=20, color=hcolor, bold=True)
        # Accent underline
        add_bar(slide, x + Inches(0.35), y + Inches(0.75), Inches(1.2), Pt(2.5), hcolor)

        yo = Inches(1.0)
        for title, desc in items:
            txt(slide, x + Inches(0.35), y + yo, w - Inches(0.7), Inches(0.35),
                title, size=13, color=WHITE, bold=True)
            yo += Inches(0.32)
            txt(slide, x + Inches(0.35), y + yo, w - Inches(0.7), Inches(0.5),
                desc, size=11, color=LIGHT_GRAY)
            yo += Inches(0.52)

    left = [
        ("Dedicated National & On-Site Team",
         "Qualified local professionals backed by international expertise \u2014 always present, always accountable."),
        ("Full Source Code Ownership",
         "No third-party dependencies. Complete freedom to customize, integrate, and evolve."),
        ("Agile Delivery Model",
         "Flexible methodology with rapid reaction time \u2014 built to adapt, not bureaucratize."),
    ]

    right = [
        ("Deep Customization Capability",
         "Features tailored precisely to customer-specific requirements at any stage."),
        ("Long-Term Strategic Partnership",
         "Shared commitment to continuous innovation and advancing data intelligence together."),
    ]

    card_col(s3, Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.3),
             "Operational Edge", BRAND_BLUE, left)
    card_col(s3, Inches(6.8), Inches(1.7), Inches(5.8), Inches(5.3),
             "Partnership Value", BRAND_ORANGE, right)

    # =========================================================
    # SLIDE 4: CLOSING / SUMMARY
    # =========================================================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s4, DARK_BG)
    top_bar(s4)

    txt(s4, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2),
        "A Framework Built to Prove Value,\nNot Just Check Boxes",
        size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER, spacing=1.15)

    add_bar(s4, Inches(5.5), Inches(3.15), Inches(2.3), Pt(3), BRAND_ORANGE)

    pillars = [
        ("Capability", "Benchmark against competitors\nwith real-world metrics", BRAND_BLUE),
        ("Flexibility", "Agile team, full source code,\ndeep customization", BRAND_ORANGE),
        ("Partnership", "Long-term commitment to\njoint innovation", ACCENT_TEAL),
    ]

    for i, (title, desc, color) in enumerate(pillars):
        x = Inches(1.2) + i * Inches(3.8)
        add_rect(s4, x, Inches(3.7), Inches(3.3), Inches(2.6), fill=CARD_BG, border=CARD_BORDER)
        add_bar(s4, x, Inches(3.7), Inches(3.3), Pt(4), color)
        txt(s4, x + Inches(0.3), Inches(4.1), Inches(2.7), Inches(0.5),
            title, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
        txt(s4, x + Inches(0.3), Inches(4.7), Inches(2.7), Inches(1.2),
            desc, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    txt(s4, Inches(0.8), Inches(6.8), Inches(5), Inches(0.4),
        "31 Concept  |  31c.io  |  Proprietary & Confidential", size=10, color=MED_GRAY)

    # Save
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
