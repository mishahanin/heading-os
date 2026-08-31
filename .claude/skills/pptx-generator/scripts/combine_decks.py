#!/usr/bin/env python3
"""Combine batch PPTX parts into one deck, keeping charts and pictures alive.

Found by the 2026-08-23 audit. `references/generation-workflow.md` prescribed
this combine loop, and the batched workflow is MANDATORY for decks over five
slides:

    for shape in slide.shapes:
        el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(el)

That copies the shape's XML and nothing else. A chart lives in its own package
part and a picture in `ppt/media/`; the shape only holds a relationship id
pointing at them. Copy the element into a slide whose part has no such
relationship and the id dangles.

Measured on python-pptx 1.0.2, combining a two-slide part holding one chart and
one picture:

    charts in package : []
    media in package  : []
    slide 2 chart   -> KeyError "no relationship with key 'rId2'"
    slide 3 picture -> KeyError "no relationship with key 'rId2'"

Zero chart parts, zero media parts. In PowerPoint those slides render blank or
broken. The cookbook ships `chart-slide.py` and `image-caption-slide.py`, so
this is the ordinary case, not an exotic one.

The fix is to walk every relationship-namespace attribute in the copied XML,
resolve it against the SOURCE part, attach the same target to the DESTINATION
part, and write back the new id.

Also preserved from the previous version: the slide background, which
`add_slide()` resets to white and shape copying does not carry.

Usage:
    python3 .claude/skills/pptx-generator/scripts/combine_decks.py \\
        --parts "$DECK_DIR/mydeck-part*.pptx" \\
        --out "$DECK_DIR/mydeck-final.pptx" \\
        --background 0E1116 \\
        --delete-parts

Guarded by tests/test_pptx_combine_keeps_charts_and_images.py.
"""
from __future__ import annotations

import argparse
import copy
import glob
import re
import shutil
import sys
from pathlib import Path

RELS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
BLANK_LAYOUT = 6

# The LAST run of digits in the stem. `(?!.*\d)` rejects any match with another
# digit after it, so `2026-08-31-deck-part12` reads 12 and not 2026.
_TRAILING_NUMBER_RE = re.compile(r"(\d+)(?!.*\d)")

# python-pptx is an optional extra (`x-heading-requires: ["documents"]`), so it
# is imported inside the functions that need it. tests/test_import_purity.py
# fails any module under the workspace that pulls a blocked optional dep at
# import time.


def hex_to_rgb(hex_color: str):
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def remap_relationships(element, src_part, dst_part) -> int:
    """Re-point every r:id in a copied shape at the destination part.

    Returns the number of relationships remapped, so a caller can assert the
    copy actually carried something rather than trusting it silently.
    """
    remapped = 0
    for node in element.iter():
        for name, rid in list(node.attrib.items()):
            if not name.startswith(f"{{{RELS_NS}}}"):
                continue
            rel = src_part.rels[rid]
            if rel.is_external:
                new_rid = dst_part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True)
            else:
                new_rid = dst_part.relate_to(rel.target_part, rel.reltype)
            node.set(name, new_rid)
            remapped += 1
    return remapped


def copy_slide(src_slide, dst_prs, background: str | None) -> int:
    """Append one slide's shapes to a new slide in dst_prs."""
    new_slide = dst_prs.slides.add_slide(dst_prs.slide_layouts[BLANK_LAYOUT])

    # add_slide() gives a default white background, and shape copying does not
    # carry the slide's own background - it is a slide property, not a shape.
    if background:
        new_slide.background.fill.solid()
        new_slide.background.fill.fore_color.rgb = hex_to_rgb(background)

    remapped = 0
    for shape in src_slide.shapes:
        element = copy.deepcopy(shape.element)
        remapped += remap_relationships(
            element, src_slide.part, new_slide.part)
        new_slide.shapes._spTree.append(element)
    return remapped


def sort_part_files(paths: list[Path]) -> list[Path]:
    """Order part files by the number in their name, not lexically.

    `references/generation-workflow.md` prescribes the UNPADDED name
    ``{name}-part1.pptx``. Batches cap at 5 slides, so any deck over 45 slides
    reaches ten parts and `sorted()` returns::

        part1, part10, part11, part2, part3, ... part9

    `combine` promises to merge "in the given order" and was handed that one, so
    a 55-slide deck shipped with slides 46-55 between slide 5 and slide 6.

    Natural sort rather than zero-padding the prescribed name: padding would fix
    only decks generated after the change, and every part file already on disk
    is unpadded. This orders both, so no transition window exists.

    A name carrying no digits is NOT dropped and does not raise -- `--parts` is a
    glob, so an operator's `deck-partial.pptx` can match it. Those sort after
    every numbered part, among themselves by name, which is stable and visible
    rather than silent.
    """
    def key(path: Path) -> tuple[int, int, str]:
        match = _TRAILING_NUMBER_RE.search(path.stem)
        if match is None:
            return (1, 0, path.name)
        return (0, int(match.group(1)), path.name)

    return sorted(paths, key=key)


def combine(part_files: list[Path], out_path: Path,
            background: str | None = None) -> Path:
    """Merge parts in the given order into out_path. Returns out_path."""
    from pptx import Presentation

    if not part_files:
        raise SystemExit("Error: no part files matched.")
    if len(part_files) == 1:
        shutil.copy(str(part_files[0]), str(out_path))
        return out_path

    combined = Presentation(str(part_files[0]))
    for part_file in part_files[1:]:
        part_prs = Presentation(str(part_file))
        for slide in part_prs.slides:
            copy_slide(slide, combined, background)
    combined.save(str(out_path))
    return out_path


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Combine batch PPTX parts into one deck.")
    parser.add_argument("--parts", required=True,
                        help="glob for the part files, e.g. '$DECK_DIR/deck-part*.pptx'")
    parser.add_argument("--out", required=True, help="path of the combined deck")
    parser.add_argument("--background",
                        help="brand background hex, e.g. 0E1116. Slides added by "
                             "the combine default to white without it.")
    parser.add_argument("--delete-parts", action="store_true",
                        help="remove the part files after a successful combine")
    args = parser.parse_args()

    part_files = sort_part_files([Path(p) for p in glob.glob(args.parts)])
    if not part_files:
        print(f"Error: no files matched {args.parts}", file=sys.stderr)
        return 2

    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    combine(part_files, out_path, args.background)

    if args.delete_parts:
        for part_file in part_files:
            if part_file.resolve() != out_path.resolve():
                part_file.unlink()

    print(f"Combined {len(part_files)} part(s) -> {out_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
