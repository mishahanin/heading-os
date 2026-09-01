#!/usr/bin/env python3
"""OPC-layer font embedding for the editable Pencil-export PPTX.

`embed_fonts()` adds the PowerPoint 'Embed fonts in the file' structures to a
saved .pptx package: a fntdata content-type default, one font part + relationship
per used typeface, and a schema-ordered <p:embeddedFontLst> with
embedTrueTypeFonts on the presentation. These tests exercise that package surgery
hermetically - no playwright, no rendering, no brand-font dependency (the embedded
bytes need not be a real font for the packaging to be correct)."""

import importlib.util
import zipfile
from pathlib import Path

from lxml import etree
import pytest
pytest.importorskip("pptx")  # F-7.1: skip on a core-only clone (needs the documents extra)

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
# hyphenated filename -> load as a module by path
_spec = importlib.util.spec_from_file_location("pencil_export", ROOT / "scripts" / "pencil-export.py")
pencil_export = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(pencil_export)

P = "http://schemas.openxmlformats.org/presentationml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT = "http://schemas.openxmlformats.org/package/2006/content-types"


def _blank_pptx(path: Path):
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))


def _fake_font(dir_: Path, name: str) -> Path:
    dir_.mkdir(parents=True, exist_ok=True)
    f = dir_ / name
    f.write_bytes(b"\x00\x01\x00\x00fake-ttf-bytes")  # not a real font; bytes are opaque to packaging
    return f


def test_embed_fonts_adds_all_opc_structures(tmp_path):
    fonts = tmp_path / "fonts"
    _fake_font(fonts, "GT Standard L Medium.ttf")
    pptx = tmp_path / "deck.pptx"
    _blank_pptx(pptx)

    n = pencil_export.embed_fonts(pptx, {"GT Standard L Medium": False}, [fonts])
    assert n == 1

    with zipfile.ZipFile(pptx) as z:
        names = set(z.namelist())
        assert "ppt/fonts/font1.fntdata" in names
        assert z.read("ppt/fonts/font1.fntdata") == b"\x00\x01\x00\x00fake-ttf-bytes"

        ctypes = etree.fromstring(z.read("[Content_Types].xml"))
        assert any(d.get("Extension") == "fntdata" and d.get("ContentType") == "application/x-fontdata"
                   for d in ctypes.findall(f"{{{CT}}}Default"))

        rels = etree.fromstring(z.read("ppt/_rels/presentation.xml.rels"))
        font_rels = [r for r in rels
                     if r.get("Type", "").endswith("/font") and r.get("Target") == "fonts/font1.fntdata"]
        assert len(font_rels) == 1
        rid = font_rels[0].get("Id")

        pres = etree.fromstring(z.read("ppt/presentation.xml"))
        assert pres.get("embedTrueTypeFonts") == "1"
        lst = pres.find(f"{{{P}}}embeddedFontLst")
        assert lst is not None
        font = lst.find(f"{{{P}}}embeddedFont/{{{P}}}font")
        assert font.get("typeface") == "GT Standard L Medium"
        reg = lst.find(f"{{{P}}}embeddedFont/{{{P}}}regular")
        assert reg.get(f"{{{R}}}id") == rid  # embeddedFontLst rId matches the relationship


def test_embed_fonts_orders_lst_after_notessz(tmp_path):
    """embeddedFontLst must sit after notesSz (CT_Presentation sequence) or
    PowerPoint rejects the file."""
    fonts = tmp_path / "fonts"
    _fake_font(fonts, "GT Standard L Medium.ttf")
    pptx = tmp_path / "deck.pptx"
    _blank_pptx(pptx)
    pencil_export.embed_fonts(pptx, {"GT Standard L Medium": False}, [fonts])

    with zipfile.ZipFile(pptx) as z:
        pres = etree.fromstring(z.read("ppt/presentation.xml"))
    children = [etree.QName(c).localname for c in pres]
    assert "embeddedFontLst" in children
    assert children.index("notesSz") < children.index("embeddedFontLst")
    # nothing after embeddedFontLst that must precede it
    assert children.index("embeddedFontLst") > children.index("sldSz")


def _fake_font_bytes(marker: bytes) -> bytes:
    return b"\x00\x01\x00\x00" + marker


def _distinct_font(dir_: Path, name: str, marker: bytes) -> Path:
    """Like `_fake_font`, but each file's bytes identify WHICH file it is.

    `_fake_font` writes the same bytes into every file, so a slot pointing at the
    wrong face is indistinguishable from one pointing at the right face. That is
    what let the oblique-intent filter go unmeasured (see the italic test below).
    """
    dir_.mkdir(parents=True, exist_ok=True)
    f = dir_ / name
    f.write_bytes(_fake_font_bytes(marker))
    return f


def test_embed_fonts_adds_italic_slot(tmp_path):
    """Both slots exist AND each carries its own face.

    Asserting only that a `<regular>` and an `<italic>` element appeared says
    nothing about which file went into either. Measured 2026-09-01 by disabling
    the oblique-intent filter in `_embeddable_font`
    (`if bool(ft & {"oblique", "italic"}) != fam_oblique: continue` -> `if False`):
    the italic slot resolved to the REGULAR face, because "GT Standard L Medium"
    ranks ahead of "GT Standard L Medium Oblique" on `len(ft - fam_core)` once the
    intent check is gone. Two slots were still written, both elements were still
    present, and this test stayed green while the deck shipped an italic run in
    an upright face. The byte markers below are what tell the two apart.
    """
    fonts = tmp_path / "fonts"
    _distinct_font(fonts, "GT Standard L Medium.ttf", b"REGULAR")
    _distinct_font(fonts, "GT Standard L Medium Oblique.ttf", b"OBLIQUE")
    pptx = tmp_path / "deck.pptx"
    _blank_pptx(pptx)

    n = pencil_export.embed_fonts(pptx, {"GT Standard L Medium": True}, [fonts])
    assert n == 2  # regular + italic

    with zipfile.ZipFile(pptx) as z:
        pres = etree.fromstring(z.read("ppt/presentation.xml"))
        rels = etree.fromstring(z.read("ppt/_rels/presentation.xml.rels"))
        target = {r.get("Id"): r.get("Target") for r in rels
                  if r.get("Type", "").endswith("/font")}
        ef = pres.find(f"{{{P}}}embeddedFontLst/{{{P}}}embeddedFont")
        assert ef.find(f"{{{P}}}regular") is not None
        assert ef.find(f"{{{P}}}italic") is not None

        got = {}
        for style in ("regular", "italic"):
            rid = ef.find(f"{{{P}}}{style}").get(f"{{{R}}}id")
            got[style] = z.read("ppt/" + target[rid])

    assert got["regular"] == _fake_font_bytes(b"REGULAR"), (
        "the regular slot does not carry the upright face")
    assert got["italic"] == _fake_font_bytes(b"OBLIQUE"), (
        "the italic slot does not carry the oblique face; the oblique-intent "
        "filter in _embeddable_font is not doing its job")


def test_embeddable_font_prefers_ttf_on_the_token_match_path_too(tmp_path):
    """The TTF preference lives on BOTH resolution paths, and both are pinned.

    `test_embeddable_font_prefers_ttf_over_otf` below uses glued filenames, which
    only the normalized-equality fallback resolves, so it pins the `.ttf` rank in
    that sort and nothing else. Measured 2026-09-01: flipping the primary
    token-match rank to prefer `.otf` left every test in this file green. The
    brand fonts ship in both formats, so the untested branch is the one the real
    export takes.
    """
    fonts = tmp_path / "fonts"
    _fake_font(fonts, "GT-Standard-L-Standard-Medium.otf")
    _fake_font(fonts, "GT-Standard-L-Standard-Medium.ttf")
    hit = pencil_export._embeddable_font("GT Standard L Medium", [fonts])
    assert hit is not None and hit.suffix == ".ttf", (
        f"token match picked {hit.name if hit else None}, not the TTF")


def test_embeddable_font_matches_glued_filename(tmp_path):
    """Typeface '31C Horizontal T03 560' must resolve the glued-name OTF
    '31CHorizontalT03-560.otf' via the normalized-equality fallback."""
    fonts = tmp_path / "fonts"
    _fake_font(fonts, "31CHorizontalT03-560.otf")
    _fake_font(fonts, "31CHorizontalT03-Clarity.otf")
    hit = pencil_export._embeddable_font("31C Horizontal T03 560", [fonts])
    assert hit is not None and hit.name == "31CHorizontalT03-560.otf"


def test_embeddable_font_token_match_respects_optical_size(tmp_path):
    """Token match must pick the Medium face, never let optical-size 'L' bleed into
    'Light'. Prefers TTF over OTF."""
    fonts = tmp_path / "fonts"
    _fake_font(fonts, "GT-Standard-L-Standard-Medium.ttf")
    _fake_font(fonts, "GT-Standard-L-Standard-Light.ttf")
    hit = pencil_export._embeddable_font("GT Standard L Medium", [fonts])
    assert hit is not None and hit.name == "GT-Standard-L-Standard-Medium.ttf"


def test_embeddable_font_prefers_ttf_over_otf(tmp_path):
    fonts = tmp_path / "fonts"
    _fake_font(fonts, "31CHorizontalT03-560.otf")
    _fake_font(fonts, "31CHorizontalT03-560.ttf")
    hit = pencil_export._embeddable_font("31C Horizontal T03 560", [fonts])
    assert hit.suffix == ".ttf"


def test_embed_fonts_noop_when_no_font_file(tmp_path):
    fonts = tmp_path / "fonts"
    fonts.mkdir()
    pptx = tmp_path / "deck.pptx"
    _blank_pptx(pptx)
    n = pencil_export.embed_fonts(pptx, {"Nonexistent Face": False}, [fonts])
    assert n == 0
    with zipfile.ZipFile(pptx) as z:
        pres = etree.fromstring(z.read("ppt/presentation.xml"))
    assert pres.find(f"{{{P}}}embeddedFontLst") is None  # unchanged
