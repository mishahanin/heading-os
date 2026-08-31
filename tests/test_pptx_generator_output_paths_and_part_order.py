#!/usr/bin/env python3
"""Two defects in `.claude/skills/pptx-generator/`, plus the one found fixing them.

**Part order (F1).** `combine_decks.py` collected its parts with
``sorted(glob.glob(...))`` while `references/generation-workflow.md` prescribes
the UNPADDED name ``{name}-part1.pptx``. Batches cap at 5 slides, so a deck over
45 slides produces ten or more parts and lexical order delivers them as::

    part1, part10, part11, part2, part3, ... part9

The function's own docstring says "Merge parts in the given order"; the order it
was handed was lexical. A 50-slide deck combined with slides 46-55 sitting
between slide 5 and slide 6.

**Output path (F7).** All 21 cookbook templates and
`generate-cookbook-preview.py` ended with::

    output = Path("title-slide.pptx")
    prs.save(output)

A bare relative path resolves against the CWD, so a run from the engine clone
drops a `.pptx` into the root of a PUBLIC repository. `config/routing-map.yaml`
has ``default: engine``, so a root-level `.pptx` routes ENGINE and
`tests/test_engine_tree_clean.py` passes it: the wall never sees the leak.
Verified 2026-08-31 -- ``find_data_artifacts(["mydeck.pptx"])`` returns ``[]``.

These files are copy-paste TEMPLATES (`SKILL.md` § Step 2 reads their first 40
lines and the agent copies the body), so the bare path is not merely a bad
default: it is the shape that propagates into every generated deck script.

The fix reuses the convention `references/generation-workflow.md` already
exports, ``$DECK_DIR``, rather than inventing a second one. Resolving through
`scripts/utils/workspace.py` instead was tested and rejected -- see
`test_the_data_root_helper_is_not_the_right_seam_for_a_template`.

**PEP-723 (found while verifying F7).** Every template carries
``#!/usr/bin/env -S uv run`` and an inline script block, so the brief assumed
they run standalone. Twenty-one of twenty-two did not: PEP 723's reference
regex is greedy over its content group, and ``# ///`` matches that group's
``# .*`` alternative, so the script block ran past its own terminator into the
second (``# /// layout``) block. Measured with uv 0.11.21::

    $ uv run .claude/skills/pptx-generator/cookbook/title-slide.py
    error: TOML parse error at line 5, column 4
      | ///

A blank line between the two blocks stops the greedy match, because the content
group requires every line to begin with ``#``.
"""
from __future__ import annotations

import ast
import importlib.util
import os
import re
import subprocess
import sys
import tomllib
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
SKILL = ROOT / ".claude" / "skills" / "pptx-generator"
COMBINER = SKILL / "scripts" / "combine_decks.py"

_spec = importlib.util.spec_from_file_location("_combine_order_under_test", COMBINER)
combine_mod = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(combine_mod)

BLANK = 6


def _template_files() -> list[Path]:
    """The 21 cookbook templates plus the preview generator."""
    return sorted(SKILL.glob("cookbook/*.py")) + \
        sorted(SKILL.glob("cookbook/carousels/*.py")) + \
        [SKILL / "generate-cookbook-preview.py"]


# Floored so a broken glob cannot make every guard below pass over an empty set.
EXPECTED_TEMPLATE_COUNT = 22


def test_the_template_corpus_is_the_size_this_file_claims():
    found = _template_files()
    assert len(found) == EXPECTED_TEMPLATE_COUNT, (
        f"expected {EXPECTED_TEMPLATE_COUNT} pptx-generator templates, found "
        f"{len(found)}: {[p.name for p in found]}. Every guard in this file "
        "iterates that glob; an empty or shrunken one asserts nothing."
    )
    for path in found:
        assert path.is_file(), f"{path} does not exist"


# ---------------------------------------------------------------------------
# F1 -- part ordering
# ---------------------------------------------------------------------------

def _make_part(path: Path, marker: str) -> None:
    """One single-slide deck whose only text is `marker`."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = marker
    p.font.size = Pt(28)
    prs.save(str(path))


def _slide_markers(path: Path) -> list[str]:
    out = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.strip())
                break
    return out


def test_a_ten_part_deck_combines_in_numeric_not_lexical_order(tmp_path):
    """The whole finding, end to end, through the real CLI.

    Eleven parts is the smallest count that separates the two orderings; a deck
    of 51-55 slides at the mandatory 5-per-batch cap produces exactly this.
    """
    for i in range(1, 12):
        _make_part(tmp_path / f"deck-part{i}.pptx", f"slide-{i:02d}")

    out = tmp_path / "deck-final.pptx"
    rc = subprocess.run(
        [sys.executable, str(COMBINER),
         "--parts", str(tmp_path / "deck-part*.pptx"),
         "--out", str(out), "--background", "0E1116"],
        capture_output=True, text=True,
    )
    assert rc.returncode == 0, rc.stderr

    assert _slide_markers(out) == [f"slide-{i:02d}" for i in range(1, 12)], (
        "parts were merged in lexical order (part1, part10, part11, part2 ...) "
        "rather than numeric order"
    )


def test_the_sort_key_orders_eleven_parts_numerically():
    names = [Path(f"/d/deck-part{i}.pptx") for i in (3, 11, 1, 10, 2, 9)]
    assert [p.name for p in combine_mod.sort_part_files(names)] == [
        "deck-part1.pptx", "deck-part2.pptx", "deck-part3.pptx",
        "deck-part9.pptx", "deck-part10.pptx", "deck-part11.pptx",
    ]


def test_lexical_sort_really_would_get_this_wrong():
    """The other jaw: if `sorted()` ever started producing numeric order, the
    test above would pass over a bug it is no longer detecting."""
    names = [f"deck-part{i}.pptx" for i in range(1, 12)]
    assert sorted(names) != names
    assert sorted(names)[1] == "deck-part10.pptx"


def test_a_part_file_with_no_trailing_number_is_kept_not_dropped_or_crashed():
    """An unexpected filename must not crash and must not vanish silently.

    `--parts` takes a glob, so an operator's `deck-partial.pptx` or a
    hand-renamed `deck-part-final.pptx` can match it.
    """
    names = [Path(p) for p in (
        "/d/deck-part2.pptx", "/d/deck-partial.pptx",
        "/d/deck-part1.pptx", "/d/deck-part-final.pptx",
    )]
    ordered = combine_mod.sort_part_files(names)
    assert len(ordered) == len(names), "a part file was dropped"
    assert {p.name for p in ordered} == {p.name for p in names}
    # Numbered parts keep their numeric order and lead; the unnumbered ones
    # follow in a stable, lexical order rather than interleaving unpredictably.
    assert [p.name for p in ordered] == [
        "deck-part1.pptx", "deck-part2.pptx",
        "deck-part-final.pptx", "deck-partial.pptx",
    ]


def test_a_multi_digit_group_is_read_as_one_number():
    names = [Path(f"/d/x-part{i}.pptx") for i in (100, 2, 20)]
    assert [p.name for p in combine_mod.sort_part_files(names)] == [
        "x-part2.pptx", "x-part20.pptx", "x-part100.pptx",
    ]


def test_a_dated_filename_reads_the_part_number_not_the_year():
    """`output.naming` is `{name}-{date}`, so a dated stem is the ordinary case.

    Added after a surviving mutation: dropping the `(?!.*\\d)` lookahead, so the
    regex takes the FIRST digit run, left every test above green. Every stem
    here then reads 2026, the numbers tie, and the tie-break falls to the
    lexical name the whole fix exists to escape.
    """
    names = [Path(f"/d/deck-2026-08-31-part{i}.pptx") for i in (2, 11, 1, 10)]
    assert [p.name for p in combine_mod.sort_part_files(names)] == [
        "deck-2026-08-31-part1.pptx", "deck-2026-08-31-part2.pptx",
        "deck-2026-08-31-part10.pptx", "deck-2026-08-31-part11.pptx",
    ]


# ---------------------------------------------------------------------------
# F7 -- output paths
# ---------------------------------------------------------------------------

DATA_SUFFIXES = {".pptx", ".pdf", ".png", ".jpg", ".jpeg", ".docx"}
SINK_CALLS = {"Path", "save", "savefig", "open", "write_bytes", "write_text"}


def bare_relative_data_sinks(tree: ast.AST) -> list[tuple[int, str]]:
    """Every data-shaped filename literal handed straight to an output sink.

    Structural, not textual: a workspace guard may not be a source grep. A
    literal joined onto a resolved directory (``OUT_DIR / "x.pptx"``) is an
    ast.BinOp operand, never a Call argument, so it is correctly not a hit.
    """
    hits: list[tuple[int, str]] = []
    for node in ast.walk(tree):
        if not isinstance(node, ast.Call):
            continue
        func = node.func
        if isinstance(func, ast.Name):
            name = func.id
        elif isinstance(func, ast.Attribute):
            name = func.attr
        else:
            continue
        if name not in SINK_CALLS:
            continue
        for arg in node.args:
            if not (isinstance(arg, ast.Constant) and isinstance(arg.value, str)):
                continue
            value = arg.value
            if "/" in value or os.sep in value:
                continue
            if Path(value).suffix.lower() in DATA_SUFFIXES:
                hits.append((arg.lineno, value))
    return hits


def resolves_an_output_root(tree: ast.AST) -> bool:
    """True when the module derives its output dir instead of assuming the CWD.

    Either from the environment (the `$DECK_DIR` convention
    `references/generation-workflow.md` already exports) or from the module's
    own location. Without this half, `Path(".") / "x.pptx"` would satisfy the
    literal check above while still writing wherever the shell happened to be.
    """
    for node in ast.walk(tree):
        if isinstance(node, ast.Attribute) and node.attr == "environ":
            return True
        if isinstance(node, ast.Name) and node.id == "__file__":
            return True
    return False


@pytest.mark.parametrize(
    "path", _template_files(), ids=lambda p: p.name)
def test_no_template_writes_to_a_bare_relative_path(path: Path):
    tree = ast.parse(path.read_text(encoding="utf-8"), filename=str(path))
    hits = bare_relative_data_sinks(tree)
    assert not hits, (
        f"{path.relative_to(ROOT)} hands a bare relative filename to an output "
        f"sink: {hits}. Run from the engine clone this writes into a PUBLIC "
        "repository, and root-level artifacts route `engine` so the "
        "engine-tree-clean wall does not catch them."
    )


@pytest.mark.parametrize(
    "path", _template_files(), ids=lambda p: p.name)
def test_every_template_resolves_its_output_root(path: Path):
    tree = ast.parse(path.read_text(encoding="utf-8"), filename=str(path))
    assert resolves_an_output_root(tree), (
        f"{path.relative_to(ROOT)} never reads os.environ or __file__, so its "
        "output directory can only be the process CWD."
    )


def test_the_detector_fires_on_the_shape_that_was_actually_there():
    """A guard with no negative case is not a guard."""
    before = ast.parse(
        'from pathlib import Path\n'
        'output = Path("title-slide.pptx")\n'
        'prs.save(output)\n'
    )
    assert bare_relative_data_sinks(before) == [(2, "title-slide.pptx")]
    assert resolves_an_output_root(before) is False


def test_the_detector_accepts_the_shape_that_replaced_it():
    after = ast.parse(
        'import os\n'
        'from pathlib import Path\n'
        'OUT_DIR = Path(os.environ["DECK_DIR"])\n'
        'output = OUT_DIR / "title-slide.pptx"\n'
        'prs.save(output)\n'
    )
    assert bare_relative_data_sinks(after) == []
    assert resolves_an_output_root(after) is True


def test_a_dot_slash_dodge_does_not_satisfy_both_halves():
    """`Path(".") / "x.pptx"` clears the literal check; the second half stops it."""
    dodge = ast.parse(
        'from pathlib import Path\n'
        'output = Path(".") / "x.pptx"\n'
    )
    assert bare_relative_data_sinks(dodge) == []
    assert resolves_an_output_root(dodge) is False


# ---------------------------------------------------------------------------
# F7, behaviour -- a template refuses rather than guessing
# ---------------------------------------------------------------------------

A_TEMPLATE = SKILL / "cookbook" / "title-slide.py"


def _runnable_copy(tmp_path: Path) -> Path:
    """A template with its `REPLACE` placeholders filled in.

    The templates ship unrunnable on purpose -- every brand colour and every
    content string is the literal `REPLACE`. `1A2B3C` is chosen because it is
    simultaneously a valid six-digit hex colour, a usable font name, and
    printable headline text, so one substitution serves all eight sites.
    """
    src = A_TEMPLATE.read_text(encoding="utf-8")
    assert '"REPLACE"' in src, "the template no longer carries placeholders"
    dst = tmp_path / "filled-title-slide.py"
    dst.write_text(src.replace('"REPLACE"', '"1A2B3C"'), encoding="utf-8")
    return dst


def test_a_template_run_without_deck_dir_refuses_and_writes_nothing(tmp_path):
    """Console-first § 4: degrade clearly, never silently.

    The guard sits at the top of `main()`, so this refuses BEFORE rendering.
    A runnable copy is used deliberately: against the raw template the process
    would exit non-zero anyway, on `int('RE', 16)`, and the test would pass
    while proving nothing about DECK_DIR.
    """
    script = _runnable_copy(tmp_path)
    cwd = tmp_path / "elsewhere"
    cwd.mkdir()
    env = {k: v for k, v in os.environ.items() if k != "DECK_DIR"}
    rc = subprocess.run(
        [sys.executable, str(script)],
        cwd=str(cwd), capture_output=True, text=True, env=env,
    )
    assert rc.returncode != 0, "a missing DECK_DIR must not be guessed at"
    assert "DECK_DIR" in (rc.stderr + rc.stdout), rc.stderr
    assert "Traceback" not in rc.stderr, "it should refuse, not crash"
    assert list(cwd.glob("*.pptx")) == [], "it wrote a deck anyway"


# The preview generator is deliberately NOT here: its output is a tracked engine
# artifact and it defaults to its own directory. Every copy-paste template is.
COPYABLE_TEMPLATES = [p for p in _template_files()
                      if p.name != "generate-cookbook-preview.py"]


@pytest.mark.parametrize(
    "template", COPYABLE_TEMPLATES, ids=lambda p: p.name)
def test_every_template_refuses_without_deck_dir(template: Path, tmp_path):
    """All 21, not just the one the render test happens to use.

    Added after a surviving mutation: replacing the refusal in `stats-slide.py`
    with ``os.environ.get("DECK_DIR", ".")`` left the whole suite green. Both
    structural checks still passed -- the module still reads `os.environ`, and
    ``Path(deck_dir) / "stats-slide.pptx"`` is still not a bare literal -- while
    the file quietly wrote to the CWD again. Only the behaviour catches it.

    The RAW template runs here, placeholders and all. That is sound precisely
    because the guard is the first thing in `main()`: it must refuse before any
    `REPLACE` colour is parsed.
    """
    env = {k: v for k, v in os.environ.items() if k != "DECK_DIR"}
    rc = subprocess.run(
        [sys.executable, str(template)],
        cwd=str(tmp_path), capture_output=True, text=True, env=env,
    )
    assert rc.returncode != 0, f"{template.name} guessed at a missing DECK_DIR"
    assert "DECK_DIR" in (rc.stderr + rc.stdout), (
        f"{template.name} failed for some other reason:\n{rc.stderr}")
    assert "Traceback" not in rc.stderr, (
        f"{template.name} crashed instead of refusing:\n{rc.stderr}")
    assert list(tmp_path.glob("*.pptx")) == [], f"{template.name} wrote anyway"


@pytest.mark.parametrize(
    "template", COPYABLE_TEMPLATES, ids=lambda p: p.name)
def test_no_template_gives_deck_dir_a_default(template: Path):
    """The structural half of the same guard, and the cheaper one.

    `os.environ.get("DECK_DIR", <anything>)` is the exact shape that survived.
    """
    tree = ast.parse(template.read_text(encoding="utf-8"))
    for node in ast.walk(tree):
        if not (isinstance(node, ast.Call)
                and isinstance(node.func, ast.Attribute)
                and node.func.attr == "get"):
            continue
        first = node.args[0] if node.args else None
        if isinstance(first, ast.Constant) and first.value == "DECK_DIR":
            assert len(node.args) == 1, (
                f"{template.name}:{node.lineno} gives DECK_DIR a fallback "
                f"default; an unset variable must refuse, not resolve."
            )


def test_a_template_run_with_deck_dir_writes_only_there(tmp_path):
    script = _runnable_copy(tmp_path)
    deck_dir = tmp_path / "decks"
    cwd = tmp_path / "elsewhere"
    deck_dir.mkdir()
    cwd.mkdir()
    env = dict(os.environ, DECK_DIR=str(deck_dir))
    rc = subprocess.run(
        [sys.executable, str(script)],
        cwd=str(cwd), capture_output=True, text=True, env=env,
    )
    assert rc.returncode == 0, rc.stderr
    assert [p.name for p in deck_dir.glob("*.pptx")] == ["title-slide.pptx"]
    assert list(cwd.glob("*.pptx")) == [], "it also wrote next to the CWD"


PREVIEW = SKILL / "generate-cookbook-preview.py"


@pytest.fixture(scope="module")
def rendered_preview(tmp_path_factory) -> Path:
    """The preview deck, rendered into a scratch dir via `$DECK_DIR`."""
    out = tmp_path_factory.mktemp("preview")
    rc = subprocess.run(
        [sys.executable, str(PREVIEW)],
        cwd=str(out), capture_output=True, text=True,
        env=dict(os.environ, DECK_DIR=str(out)),
    )
    assert rc.returncode == 0, rc.stderr
    deck = out / "cookbook-preview.pptx"
    assert deck.is_file(), f"nothing rendered; stdout={rc.stdout}"
    return deck


def test_the_preview_honours_deck_dir_and_leaves_the_skill_tree_alone(
        rendered_preview):
    """It defaults to its own directory, but `$DECK_DIR` must win.

    Without that, a test render would overwrite the tracked
    `.claude/skills/pptx-generator/cookbook-preview.pptx`.
    """
    assert rendered_preview.parent != SKILL


def test_the_floating_cards_preview_renders_its_card_descriptions(
        rendered_preview):
    """`desc` was unpacked from `cards` and never drawn.

    The preview exists to show what each layout looks like, and
    `cookbook/floating-cards-slide.py` ships a `card_desc` textbox. The preview
    advertised the layout as title-only, so a planning agent reading it would
    size content for a card that has no body.
    """
    texts = {
        shape.text_frame.text.strip()
        for slide in Presentation(str(rendered_preview)).slides
        for shape in slide.shapes
        if shape.has_text_frame
    }
    assert "Context" in texts, "the floating-cards slide is missing entirely"
    for desc in ("Feed the model what it needs",
                 "Give it ability to act",
                 "Refine until done"):
        assert desc in texts, f"card description not rendered: {desc!r}"


def test_the_circular_hero_preview_labels_its_orbit_nodes(rendered_preview):
    """The same omission, found by scanning for the same shape.

    `item` was bound by the loop and only ever used to place a 0.16in dot, so
    the preview drew six anonymous dots. `cookbook/circular-hero-slide.py:193`
    sets `p.text = item`.
    """
    texts = {
        shape.text_frame.text.strip()
        for slide in Presentation(str(rendered_preview)).slides
        for shape in slide.shapes
        if shape.has_text_frame
    }
    for label in ("Docs", "Examples", "Constraints",
                  "History", "Tools", "Goals"):
        assert label in texts, f"orbit node not labelled: {label!r}"


def test_no_preview_builder_binds_a_value_it_never_draws(rendered_preview):
    """The guard that found the second instance, kept as the guard.

    Both defects had one shape: a name bound by a `for` target and never read.
    Neither is caught by any renderer test that does not know to look for the
    missing string, so this asserts the shape itself across every builder.
    """
    tree = ast.parse(PREVIEW.read_text(encoding="utf-8"))
    offenders = {}
    for fn in (n for n in ast.walk(tree) if isinstance(n, ast.FunctionDef)):
        bound = {
            t.id
            for node in ast.walk(fn) if isinstance(node, ast.For)
            for t in ast.walk(node.target) if isinstance(t, ast.Name)
        }
        loaded = {
            n.id for n in ast.walk(fn)
            if isinstance(n, ast.Name) and isinstance(n.ctx, ast.Load)
        }
        dead = sorted(bound - loaded - {"_"})
        if dead:
            offenders[fn.name] = dead
    assert not offenders, (
        f"loop variables bound but never rendered: {offenders}. Each is content "
        "the preview claims to show and does not."
    )


def test_the_data_root_helper_is_not_the_right_seam_for_a_template():
    """Why `$DECK_DIR` and not `get_outputs_dir()`.

    The helper IS importable from a PEP-723 standalone run (measured
    2026-08-31, uv 0.11.21: an absolute `sys.path.insert` plus
    `from scripts.utils.workspace import get_outputs_dir` resolved fine inside
    the ephemeral env). It is still the wrong seam here, for a reason that is
    about these files being templates rather than about importability:
    `get_data_root()`'s last rule answers `<workspace_root>/examples`, so on a
    clone with no data overlay -- every public clone -- the helper points back
    INSIDE the engine tree. A template must also survive being copied somewhere
    else, which a `__file__`-relative walk up to the engine root does not.
    """
    from scripts.utils.workspace import get_data_root  # noqa: F401
    import inspect
    src = inspect.getsource(get_data_root)
    assert "examples" in src, (
        "get_data_root() no longer falls back into the engine tree; re-evaluate "
        "the choice this test records."
    )


# ---------------------------------------------------------------------------
# PEP 723 -- the inline script block must actually parse
# ---------------------------------------------------------------------------

# Verbatim from PEP 723's reference implementation.
PEP723_RE = (
    r'(?m)^# /// (?P<type>[a-zA-Z0-9-]+)$\s(?P<content>(^#(| .*)$\s)+)^# ///$'
)


def pep723_script_block(source: str) -> str | None:
    for match in re.finditer(PEP723_RE, source):
        if match.group("type") != "script":
            continue
        return "".join(
            line[2:] if len(line) > 1 else line[1:]
            for line in match.group("content").splitlines(keepends=True)
        )
    return None


@pytest.mark.parametrize(
    "path", _template_files(), ids=lambda p: p.name)
def test_every_template_has_a_parseable_inline_script_block(path: Path):
    """21 of 22 failed this before 2026-08-31.

    The shebang is `uv run`, so an unparseable block means the file cannot be
    executed the way it advertises. uv reported:
    `TOML parse error at line 5, column 4 | ///`.
    """
    content = pep723_script_block(path.read_text(encoding="utf-8"))
    assert content is not None, f"{path.name} has no PEP-723 script block"
    meta = tomllib.loads(content)
    assert meta["dependencies"] == ["python-pptx==1.0.2"], (
        f"{path.name} declares {meta.get('dependencies')!r}; the inline pin must "
        "match the `python-pptx==1.0.2` pin in pyproject.toml"
    )


def test_the_greedy_regex_really_does_swallow_an_adjacent_block():
    """The failure mode, kept on file. Without the blank line, the script
    block's content runs past its own terminator into the `layout` block."""
    glued = (
        "# /// script\n"
        '# dependencies = ["python-pptx==1.0.2"]\n'
        "# ///\n"
        "# /// layout\n"
        '# name = "x"\n'
        "# ///\n"
    )
    content = pep723_script_block(glued)
    assert content is not None
    assert "///" in content, "the terminator was not swallowed; premise gone"
    with pytest.raises(tomllib.TOMLDecodeError):
        tomllib.loads(content)

    separated = glued.replace("# ///\n# /// layout", "# ///\n\n# /// layout")
    assert "///" not in pep723_script_block(separated)
