"""Combining PPTX batches must not strip charts and pictures.

Found by the 2026-08-23 audit. `references/generation-workflow.md` prescribed a
combine loop that deep-copies each shape's XML into a new slide::

    el = copy.deepcopy(shape.element)
    new_slide.shapes._spTree.append(el)

A chart lives in its own package part and a picture in `ppt/media/`; the shape
element holds only a relationship id pointing at them. Copied into a slide whose
part carries no such relationship, the id dangles.

Measured on python-pptx 1.0.2 before the fix, combining a part with one chart
slide and one picture slide::

    charts in package : []
    media in package  : []
    slide 2 chart   -> KeyError "no relationship with key 'rId2'"
    slide 3 picture -> KeyError "no relationship with key 'rId2'"

The batched workflow is MANDATORY for decks over five slides and the cookbook
ships `chart-slide.py` and `image-caption-slide.py`, so this was the ordinary
path, not an edge case. `test_the_old_naive_loop_still_reproduces_the_bug` keeps
the original failure on file: if it ever stops reproducing, this whole guard is
measuring nothing.
"""
from __future__ import annotations

import copy
import importlib.util
import zipfile
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
PIL = pytest.importorskip("PIL.Image")

from pptx import Presentation  # noqa: E402
from pptx.chart.data import CategoryChartData  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE  # noqa: E402
from pptx.util import Inches  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
COMBINER = (ROOT / ".claude" / "skills" / "pptx-generator" / "scripts"
            / "combine_decks.py")

_spec = importlib.util.spec_from_file_location("_combine_under_test", COMBINER)
combine_mod = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(combine_mod)

BLANK = 6


@pytest.fixture
def parts(tmp_path: Path) -> list[Path]:
    """Two part files: plain text, then one chart slide and one picture slide."""
    png = tmp_path / "pic.png"
    PIL.new("RGB", (80, 60), (200, 30, 30)).save(png)

    p1 = Presentation()
    s = p1.slides.add_slide(p1.slide_layouts[BLANK])
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "part one"
    part1 = tmp_path / "deck-part1.pptx"
    p1.save(part1)

    p2 = Presentation()
    s = p2.slides.add_slide(p2.slide_layouts[BLANK])
    data = CategoryChartData()
    data.categories = ["a", "b"]
    data.add_series("revenue", (1.0, 2.0))
    s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                       Inches(1), Inches(1), Inches(4), Inches(3), data)
    s2 = p2.slides.add_slide(p2.slide_layouts[BLANK])
    s2.shapes.add_picture(str(png), Inches(1), Inches(1), Inches(2))
    part2 = tmp_path / "deck-part2.pptx"
    p2.save(part2)

    return [part1, part2]


def _package(path: Path) -> list[str]:
    return zipfile.ZipFile(path).namelist()


# ------------------------------------------------------------ the regression

def test_the_old_naive_loop_still_reproduces_the_bug(parts, tmp_path: Path):
    """Pins the defect. Without this, the guard below could pass vacuously."""
    combined = Presentation(str(parts[0]))
    src = Presentation(str(parts[1]))
    for slide in src.slides:
        new_slide = combined.slides.add_slide(combined.slide_layouts[BLANK])
        for shape in slide.shapes:
            new_slide.shapes._spTree.append(copy.deepcopy(shape.element))
    naive = tmp_path / "naive.pptx"
    combined.save(naive)

    names = _package(naive)
    assert [n for n in names if "charts/chart" in n] == []
    assert [n for n in names if "media/" in n] == []

    out = Presentation(str(naive))
    slides = list(out.slides)
    with pytest.raises(KeyError):
        _ = slides[1].shapes[0].chart.plots
    with pytest.raises(KeyError):
        _ = slides[2].shapes[0].image.blob


# ------------------------------------------------------------------ the fix

def test_the_combiner_carries_the_chart_part(parts, tmp_path: Path):
    out = tmp_path / "final.pptx"
    combine_mod.combine(parts, out, background="0E1116")
    assert [n for n in _package(out) if "charts/chart" in n]

    chart = list(Presentation(str(out)).slides)[1].shapes[0].chart
    assert len(chart.plots) == 1
    series = chart.plots[0].series[0]
    assert series.name == "revenue"
    assert list(series.values) == [1.0, 2.0]


def test_the_combiner_carries_the_charts_embedded_workbook(parts, tmp_path: Path):
    """A chart's data lives in a second part behind a second relationship."""
    out = tmp_path / "final.pptx"
    combine_mod.combine(parts, out, background="0E1116")
    assert [n for n in _package(out) if "embeddings/" in n]


def test_the_combiner_carries_the_image_bytes(parts, tmp_path: Path):
    out = tmp_path / "final.pptx"
    combine_mod.combine(parts, out, background="0E1116")
    assert [n for n in _package(out) if "media/" in n]

    picture = list(Presentation(str(out)).slides)[2].shapes[0]
    assert len(picture.image.blob) > 0


def test_the_combiner_keeps_the_slide_background(parts, tmp_path: Path):
    """The bug this section already documented must stay fixed."""
    out = tmp_path / "final.pptx"
    combine_mod.combine(parts, out, background="0E1116")
    for slide in list(Presentation(str(out)).slides)[1:]:
        rgb = slide.background.fill.fore_color.rgb
        assert str(rgb) == "0E1116", str(rgb)


def test_slide_and_shape_counts_are_preserved(parts, tmp_path: Path):
    out = tmp_path / "final.pptx"
    combine_mod.combine(parts, out, background="0E1116")
    slides = list(Presentation(str(out)).slides)
    assert len(slides) == 3
    assert [len(s.shapes) for s in slides] == [1, 1, 1]


def test_a_single_part_is_copied_straight_through(parts, tmp_path: Path):
    out = tmp_path / "final.pptx"
    combine_mod.combine([parts[0]], out)
    assert out.exists()
    assert len(list(Presentation(str(out)).slides)) == 1


def test_remap_counts_what_it_moved(parts, tmp_path: Path):
    """A remapper that silently did nothing would pass a shallow check."""
    dst = Presentation(str(parts[0]))
    src = Presentation(str(parts[1]))
    moved = sum(combine_mod.copy_slide(s, dst, "0E1116") for s in src.slides)
    assert moved >= 2, f"only remapped {moved} relationship(s)"


# ------------------------------------------------------- the docs match the code

def test_the_workflow_reference_points_at_the_script():
    doc = (ROOT / ".claude" / "skills" / "pptx-generator" / "references"
           / "generation-workflow.md").read_text(encoding="utf-8")
    assert "combine_decks.py" in doc
    assert "new_slide.shapes._spTree.append(el)\n```" not in doc, (
        "the reference still prescribes the naive loop as the thing to run"
    )
