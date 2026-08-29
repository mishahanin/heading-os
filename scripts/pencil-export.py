#!/usr/bin/env python3
"""Render a Pencil deck to PNG / PDF / PPTX / self-contained HTML.

The Pencil MCP `export_nodes` tool (native PNG/PDF) is broken on WSL: its bundled
path translator prepends `\\\\wsl.localhost\\<distro>\\` inconsistently, so it cannot
write per-slide images from a WSL-resident or C:\\ .pen. `export_html` in the same
MCP resolves paths correctly, so the reliable pipeline is:

    1. (agent, MCP)  export_html(filePath=<active .pen>, outputPath=<deck>/pencil/deck.html,
                                 nodeIds=[frames in deck order], format="html-css")
       -> one HTML with every frame stacked, images referenced relative to the .pen dir.
    2. (this script) resolve/inject brand fonts, render each 1920x1080 frame in
       isolation via chromium (so overlapping absolutely-positioned frames cannot
       bleed), then assemble PDF + PPTX and, optionally, a portable self-contained HTML.

This script owns step 2. It never calls Pencil; run it on the HTML the agent
produced. Fonts are supplied via one or more --fonts-dir paths (kept out of the
engine so no brand assets are hardcoded here).

PPTX defaults to the EDITABLE twin, `<stem>.pptx`: a visually identical deck where
each slide is the text-less Pencil render as a full-bleed background image with
native, editable PowerPoint text boxes laid on top at the same coordinates and
style (branding/graphics stay baked in the background). The brand fonts used on the
runs are embedded into the file (via --fonts-dir) so it renders identically on a
machine without them installed - editability is the whole point of a PPTX.

The image-per-slide "locked look" PPTX is opt-in via the `pptx-flat` format: it is
byte-frozen and portable (needs no fonts) but not editable, and is written as
`<stem> (ready to be shared with the world).pptx`. PDF is always image-per-slide.

Usage:
    python scripts/pencil-export.py \\
      --html  <deck>/pencil/deck.html \\
      --out-dir <deck>/export \\
      --fonts-dir /path/to/brand/fonts \\
      [--stem deck] [--formats png,pdf,html,pptx,pptx-flat] \\
      [--keep-in-bg Footer --keep-in-bg Logo] \\
      [--width 1920] [--height 1080] [--scale 2] \\
      [--image-format jpeg --quality 82] [--verbose]

Requires: playwright (chromium) and python-pptx in the active venv.
"""

import argparse
import base64
import glob
import mimetypes
import os
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from scripts.utils.colors import GREEN, YELLOW, RED, CYAN, GRAY, BOLD, RESET

SLIDE_SEL = "div[data-pencil-name^='Slide-']"
FONT_EXT = (".woff2", ".woff", ".otf", ".ttf")

# Editable mode: per-slide Pencil components whose text is branding/decoration,
# not editable content. Their text stays baked in the background image; a text
# node is kept in the background when it (or an ancestor) carries one of these
# data-pencil-name values. Extend per deck with --keep-in-bg.
DECOR_DEFAULT = ("Footer", "Logo", "OrangeCorner", "Watermark", "WatermarkDark",
                 "Icon", "Check", "CheckX", "NumBadge", "LogoMark", "LogoBy")

# Runs in the browser over the Pencil export HTML. For every slide frame it
# collects each content-text BLOCK's box + style (in css px, relative to the
# frame) and marks it data-ov so the background pass can hide it. Decorative text
# (closest to a --keep-in-bg name) and icon-font glyphs are skipped.
#
# "Block", not "leaf". The rule used to be `childElementCount > 0 -> skip`, so a
# paragraph carrying ANY inline markup was skipped: `<p>Our <strong>2026</strong>
# plan</p>` has one element child, so the `<p>` was never extracted, never marked
# `data-ov`, and never hidden -- its text stayed baked into a background the tool
# calls "text-less", while `<strong>` alone became a floating text box. The words
# around the bold ran silently missing from the editable deck.
#
# An element qualifies when it holds text and every text-bearing descendant of it
# is INLINE, i.e. it is one visual paragraph. Its descendants are then marked
# covered so nothing is extracted twice; hiding the parent hides them anyway.
EXTRACT_JS = r"""
(decor) => {
  const skipSel = decor.map(n => `[data-pencil-name="${n}"]`).join(',');
  const slides = [...document.querySelectorAll("[data-pencil-name^='Slide-']")];
  const out = [];
  const isInline = (e) => {
    const d = getComputedStyle(e).display;
    return d === 'inline' || d === 'inline-block' || d === 'contents';
  };
  const hasBlockText = (el) => [...el.querySelectorAll('*')].some(
    (d) => d.textContent && d.textContent.trim() && !isInline(d));
  for (const s of slides) {
    const sr = s.getBoundingClientRect();
    const items = [];
    for (const el of s.querySelectorAll('*')) {
      if (el.hasAttribute('data-ov-in')) continue;            // inside a box already taken
      const txt = el.textContent;
      if (!txt || !txt.trim()) continue;
      if (hasBlockText(el)) continue;                         // a container, not a paragraph
      if (skipSel && el.closest(skipSel)) continue;           // decorative -> stays in bg
      const cs = getComputedStyle(el);
      const fs = parseFloat(cs.fontSize);
      if (!fs) continue;
      const fam = cs.fontFamily.split(',')[0].replace(/["']/g, '').trim();
      if (/material|lucide|feather|phosphor/i.test(fam)) continue;  // icon fonts
      const r = el.getBoundingClientRect();
      if (r.width < 1 || r.height < 1) continue;
      let lh = cs.lineHeight, lhr = 1.2;
      if (lh && lh.endsWith('px')) lhr = parseFloat(lh) / fs;
      items.push({text: txt, x: r.left - sr.left, y: r.top - sr.top,
                  w: r.width, h: r.height, fs: fs, fam: fam,
                  italic: /italic|oblique/i.test(cs.fontStyle),
                  color: cs.color, align: cs.textAlign, lhr: lhr,
                  ws: cs.whiteSpace});
      el.setAttribute('data-ov', '1');
      for (const d of el.querySelectorAll('*')) d.setAttribute('data-ov-in', '1');
    }
    // What is STILL baked in: a text node covered by no extracted box and not
    // decorative. Counted rather than assumed away -- the caller used to print
    // "text-less backgrounds", which is a claim about the whole frame that this
    // pass never established.
    let leftover = 0;
    const w = document.createTreeWalker(s, NodeFilter.SHOW_TEXT);
    while (w.nextNode()) {
      const n = w.currentNode;
      if (!n.nodeValue || !n.nodeValue.trim()) continue;
      const p = n.parentElement;
      if (!p || p.closest('[data-ov]')) continue;
      if (skipSel && p.closest(skipSel)) continue;
      leftover++;
    }
    out.push({name: s.getAttribute('data-pencil-name'),
              id: s.getAttribute('data-pencil-id'), items: items,
              leftover: leftover});
  }
  return out;
}
"""


def _slide_files(png_dir: Path):
    """Ordered slide renders, PNG or JPG."""
    return sorted(list(png_dir.glob("slide-*.png")) + list(png_dir.glob("slide-*.jpg")))


# ============================================================
# Fonts
# ============================================================

def _norm(s: str) -> str:
    """Lowercase, keep alphanumerics only, so 'GT Standard L Medium' matches
    'GT-Standard-L-Standard-Medium'."""
    return re.sub(r"[^a-z0-9]", "", s.lower())


def _find_font_file(family: str, font_dirs, by_basename=None):
    """Resolve a font-family (or exact basename) to a file under the font dirs.

    Matching rule: every token of the family must appear in the file's normalized
    name, and the family's oblique/italic intent must match the file's (so
    'GT Standard L Medium' does not grab the '-Medium-Oblique' face)."""
    for d in font_dirs:
        if by_basename:
            hit = d / by_basename
            if hit.exists():
                return hit
    if by_basename:
        # basename given but not found verbatim; fall through to fuzzy on the stem
        family = Path(by_basename).stem

    fam_tokens = [t for t in re.split(r"[^a-z0-9]+", family.lower()) if t]
    fam_oblique = any(t in ("oblique", "italic") for t in fam_tokens)
    fam_core = {t for t in fam_tokens if t not in ("oblique", "italic")}

    best = None
    for d in font_dirs:
        if not d.exists():
            continue
        for f in d.rglob("*"):
            if f.suffix.lower() not in FONT_EXT:
                continue
            # exact-token match, not substring: a single-letter optical-size
            # token (L/M/S) must equal a filename token, never match inside
            # another word (e.g. 'l' inside 'light').
            file_tokens = {t for t in re.split(r"[^a-z0-9]+", f.stem.lower()) if t}
            if not fam_core.issubset(file_tokens):
                continue
            file_oblique = bool(file_tokens & {"oblique", "italic"})
            if file_oblique != fam_oblique:
                continue
            # prefer woff2, then fewest extra tokens, then shortest name
            rank = (0 if f.suffix.lower() == ".woff2" else 1,
                    len(file_tokens - fam_core), len(f.stem))
            if best is None or rank < best[0]:
                best = (rank, f)
    return best[1] if best else None


def _data_uri(path: Path) -> str:
    mime = mimetypes.guess_type(str(path))[0]
    if not mime:
        mime = {
            ".woff2": "font/woff2", ".woff": "font/woff", ".otf": "font/otf",
            ".ttf": "font/ttf",
        }.get(path.suffix.lower(), "application/octet-stream")
    return f"data:{mime};base64," + base64.b64encode(path.read_bytes()).decode()


def resolve_fonts(html: str, font_dirs, verbose=False):
    """Return HTML with every referenced font-family backed by an embedded
    @font-face (base64). Existing @font-face rules that point at missing
    `assets/*` files are re-pointed at resolved fonts too."""
    # `font-family:'Brand'` is valid CSS and Pencil can emit it. Matching only
    # the double-quoted form made such a family invisible to BOTH sets at once:
    # never "used", so step 2 synthesised no @font-face, and never "declared",
    # so the verbose "no font file for used family" line could not fire either.
    # The font fell back to a system default with nothing printed anywhere.
    used = set(re.findall(r"""font-family:\s*["']([^"']+)["']""", html))
    declared = set(re.findall(
        r"""@font-face\s*\{[^}]*?font-family:\s*["']([^"']+)["']""", html, re.S,
    ))

    # 1. re-embed srcs of already-declared @font-face rules (assets/* won't exist)
    def _embed_src(m):
        fam, url = m.group("fam"), m.group("url")
        if url.startswith("data:"):
            return m.group(0)
        f = _find_font_file(fam, font_dirs, by_basename=Path(url).name)
        if not f:
            if verbose:
                print(f"{YELLOW}  font not found for declared @font-face {fam} ({url}){RESET}")
            return m.group(0)
        # Splice over the `url` group's own span rather than string-replacing
        # the quoted spellings. The pattern deliberately accepts an UNQUOTED
        # `src:url(assets/brand.woff2)`, but the replacement only knew `"url"`
        # and `'url'`, so on the unquoted form both calls no-opped, the rule
        # kept pointing at the `assets/*` file that does not exist in the
        # export, and the brand font fell back silently - no warning, because
        # the font file itself WAS found. Measured 2026-08-30.
        whole = m.group(0)
        start = m.start("url") - m.start()
        end = m.end("url") - m.start()
        return whole[:start] + _data_uri(f) + whole[end:]

    html = re.sub(
        r"""@font-face\s*\{[^}]*?font-family:\s*["'](?P<fam>[^"']+)["']"""
        r"""[^}]*?src:\s*url\(["']?(?P<url>[^"')]+)["']?\)[^}]*?\}""",
        _embed_src, html, flags=re.S,
    )

    # 2. synthesize @font-face for families used but never declared
    faces = []
    for fam in sorted(used - declared):
        f = _find_font_file(fam, font_dirs)
        if not f:
            if verbose:
                print(f"{YELLOW}  no font file for used family {fam}{RESET}")
            continue
        faces.append(
            "@font-face{font-family:\"%s\";src:url(\"%s\") format(\"%s\");"
            "font-display:swap;font-weight:400;}"
            % (fam, _data_uri(f), "woff2" if f.suffix.lower() == ".woff2" else "opentype")
        )
        if verbose:
            print(f"{GRAY}  + @font-face {fam} <- {f.name}{RESET}")

    if faces:
        block = "<style>" + "".join(faces) + "</style>"
        html = html.replace("<head>", "<head>" + block, 1) if "<head>" in html else block + html
    return html


# ============================================================
# Render
# ============================================================

def render_pngs(work_html: Path, png_dir: Path, width, height, scale,
                img_format="png", quality=85, verbose=False):
    """Screenshot every slide frame in isolation to png_dir/slide-NN.<ext>.

    img_format 'jpeg' (with quality) yields far smaller files than lossless PNG
    - the same win NXPowerLite gives on the PPTX, but baked in. Chromium encodes
    JPEG natively (the Pillow JPEG encoder is unavailable in this venv)."""
    from playwright.sync_api import sync_playwright

    ext = "jpg" if img_format == "jpeg" else "png"
    png_dir.mkdir(parents=True, exist_ok=True)
    for old in list(png_dir.glob("slide-*.png")) + list(png_dir.glob("slide-*.jpg")):
        old.unlink()

    count = 0
    with sync_playwright() as p:
        browser = p.chromium.launch(args=["--force-color-profile=srgb"])
        page = browser.new_page(viewport={"width": width, "height": height},
                                device_scale_factor=scale)
        page.goto(work_html.as_uri())
        page.wait_for_load_state("networkidle")
        page.evaluate("async () => { await document.fonts.ready; }")
        page.evaluate(
            "async () => { const i=[...document.images];"
            "await Promise.all(i.map(x=>x.complete?1:new Promise(r=>{x.onload=x.onerror=r;}))); }"
        )
        frames = page.query_selector_all(SLIDE_SEL)
        if not frames:
            raise SystemExit(f"{RED}No slide frames matched {SLIDE_SEL} in {work_html}{RESET}")
        # report frames sharing a top (Pencil canvas overlap)
        boxes = page.evaluate(
            "(sel)=>[...document.querySelectorAll(sel)].map(f=>({n:f.getAttribute('data-pencil-name'),"
            "t:Math.round(f.getBoundingClientRect().top+window.scrollY)}))", SLIDE_SEL,
        )
        tops = {}
        for b in boxes:
            tops.setdefault(b["t"], []).append(b["n"])
        for t, names in tops.items():
            if len(names) > 1:
                print(f"{YELLOW}  overlap at top {t}: {names} (isolation handles it){RESET}")
        for i, fr in enumerate(frames, 1):
            page.evaluate(
                "([sel,idx])=>{const fs=[...document.querySelectorAll(sel)];"
                "fs.forEach((f,j)=>{f.style.visibility=(j===idx)?'visible':'hidden';});}",
                [SLIDE_SEL, i - 1],
            )
            out = png_dir / f"slide-{i:02d}.{ext}"
            if img_format == "jpeg":
                fr.screenshot(path=str(out), type="jpeg", quality=quality)
            else:
                fr.screenshot(path=str(out))
            count += 1
        browser.close()
    if verbose:
        q = f" q{quality}" if img_format == "jpeg" else ""
        print(f"{GRAY}  rendered {count} slides at {width*scale}x{height*scale} ({img_format}{q}){RESET}")
    return count


def build_pdf(png_dir: Path, pdf_path: Path, width, height):
    """Exact one-image-per-page 16:9 PDF via chromium print (no libjpeg needed)."""
    from playwright.sync_api import sync_playwright

    pngs = _slide_files(png_dir)
    pages = "\n".join(f'<div class="pg"><img src="{p.as_uri()}"></div>' for p in pngs)
    html = (
        "<!doctype html><html><head><meta charset='utf-8'><style>"
        f"@page{{size:{width}px {height}px;margin:0;}}html,body{{margin:0;padding:0;}}"
        f".pg{{width:{width}px;height:{height}px;overflow:hidden;page-break-after:always;break-after:page;}}"
        ".pg:last-child{page-break-after:auto;break-after:auto;}"
        f"img{{display:block;width:{width}px;height:{height}px;}}</style></head><body>{pages}</body></html>"
    )
    tmp = pdf_path.parent / "_print.html"
    tmp.write_text(html, encoding="utf-8")
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page()
        page.goto(tmp.as_uri())
        page.wait_for_load_state("networkidle")
        page.evaluate(
            "async () => { const i=[...document.images];"
            "await Promise.all(i.map(x=>x.complete?1:new Promise(r=>{x.onload=x.onerror=r;}))); }"
        )
        page.pdf(path=str(pdf_path), width=f"{width}px", height=f"{height}px",
                 print_background=True,
                 margin={"top": "0", "bottom": "0", "left": "0", "right": "0"})
        browser.close()
    tmp.unlink()
    return len(pngs)


def build_pptx(png_dir: Path, pptx_path: Path):
    """Full-bleed 16:9 PPTX, one picture per slide."""
    from pptx import Presentation
    from pptx.util import Inches

    pngs = _slide_files(png_dir)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in pngs:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(pptx_path))
    return len(pngs)


# ============================================================
# Editable hybrid (text-less background + native text boxes)
# ============================================================

def render_editable(work_html: Path, bg_dir: Path, width, height, scale,
                    decor_names, img_format="jpeg", quality=82, verbose=False):
    """Extract every content-text box, then render each slide's background with
    that content text hidden (and all other slide frames hidden, so overlapping
    canvas frames cannot bleed). Returns the ordered per-slide extraction data;
    writes text-less backgrounds to bg_dir/slide-NN.<ext>."""
    from playwright.sync_api import sync_playwright

    ext = "jpg" if img_format == "jpeg" else "png"
    bg_dir.mkdir(parents=True, exist_ok=True)
    for old in list(bg_dir.glob("slide-*.png")) + list(bg_dir.glob("slide-*.jpg")):
        old.unlink()

    with sync_playwright() as p:
        browser = p.chromium.launch(args=["--force-color-profile=srgb"])
        page = browser.new_page(viewport={"width": width, "height": height},
                                device_scale_factor=scale)
        page.goto(work_html.as_uri())
        page.wait_for_load_state("networkidle")
        page.evaluate("async () => { await document.fonts.ready; }")
        page.evaluate(
            "async () => { const i=[...document.images];"
            "await Promise.all(i.map(x=>x.complete?1:new Promise(r=>{x.onload=x.onerror=r;}))); }"
        )
        data = page.evaluate(EXTRACT_JS, list(decor_names))
        if not data:
            raise SystemExit(f"{RED}No slide frames matched for editable export{RESET}")
        # hide content text (its own rule wins even when the frame is shown)
        page.add_style_tag(content="[data-ov]{visibility:hidden !important}")
        page.wait_for_timeout(150)
        for i, d in enumerate(data, 1):
            page.evaluate(
                "(cur)=>{document.querySelectorAll(\"[data-pencil-name^='Slide-']\")"
                ".forEach(e=>{e.style.visibility=(e.getAttribute('data-pencil-id')===cur)?'':'hidden';});}",
                d["id"])
            # `data-pencil-id` can be absent, in which case getAttribute gave
            # None, the selector looked for the literal string "None", and the
            # AttributeError from `fr.screenshot` aborted the whole export after
            # a full render pass with nothing that named the offending slide.
            if not d["id"]:
                raise SystemExit(
                    f"pencil-export: slide {i} ({d.get('name') or 'unnamed'}) has no "
                    f"data-pencil-id; cannot render the editable export")
            fr = page.query_selector(f'[data-pencil-id="{d["id"]}"]')
            if fr is None:
                raise SystemExit(
                    f"pencil-export: no frame matches data-pencil-id={d['id']!r} "
                    f"for slide {i}")
            out = bg_dir / f"slide-{i:02d}.{ext}"
            if img_format == "jpeg":
                fr.screenshot(path=str(out), type="jpeg", quality=quality)
            else:
                fr.screenshot(path=str(out))
        browser.close()
    message, incomplete = extraction_summary(data)
    if incomplete:
        print(message, file=sys.stderr)
    elif verbose:
        print(message)
    return data


def extraction_summary(data) -> tuple[str, bool]:
    """One line describing what the extraction ESTABLISHED, plus whether it was
    incomplete.

    The old line read "N text-less backgrounds", which is a claim about the
    whole frame. The method only hides the boxes it managed to extract, so a
    paragraph it skipped stayed visible in the image and that sentence covered
    for it -- exactly the silent partial extraction the leaf-node rule caused.
    `leftover` is the measurement that makes the difference sayable.
    """
    nb = sum(len(d["items"]) for d in data)
    left = sum(d.get("leftover", 0) for d in data)
    if left:
        return (f"{YELLOW}  {len(data)} backgrounds, {nb} content text boxes; "
                f"{left} text node(s) could NOT be extracted and stay baked "
                f"into the background (decorative text excluded).{RESET}", True)
    return (f"{GRAY}  {len(data)} backgrounds with no content text left in "
            f"them, {nb} content text boxes{RESET}", False)


def _rgb(s: str):
    from pptx.dml.color import RGBColor
    m = re.findall(r"[\d.]+", s or "")
    if len(m) >= 3:
        return RGBColor(int(round(float(m[0]))), int(round(float(m[1]))), int(round(float(m[2]))))
    return RGBColor(0, 0, 0)


def build_editable_pptx(bg_dir: Path, data, pptx_path: Path, width, height,
                        font_dirs=None, verbose=False):
    """Hybrid PPTX: each slide is its text-less background image full-bleed, with
    native (editable) text boxes laid on top at the extracted coordinates and
    style. Coordinate map: EMU_per_px = 12192000/width, font pt = px*EMU_per_px/12700
    (px == 1/144 in on a 1920px 16:9 slide). Single-line boxes never wrap (a wider
    renderer would otherwise break the last word); multi-line boxes keep the
    extracted line-height. When font_dirs are given, the brand typefaces used on the
    runs are embedded into the package so the file renders identically without them
    installed. Returns (slide_count, embedded_font_count)."""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
                 "start": PP_ALIGN.LEFT, "end": PP_ALIGN.RIGHT}
    bgs = _slide_files(bg_dir)
    if len(bgs) != len(data):
        raise SystemExit(f"{RED}editable: {len(bgs)} backgrounds != {len(data)} slides{RESET}")

    emu_px = 12192000 / width
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(round(12192000 * height / width))
    blank = prs.slide_layouts[6]
    fam_italic = {}  # typeface -> True if any run in it is italic (for embedding)
    for d, bg in zip(data, bgs):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(bg), 0, 0, width=prs.slide_width, height=prs.slide_height)
        for it in d["items"]:
            fam_italic[it["fam"]] = fam_italic.get(it["fam"], False) or bool(it["italic"])
            tb = slide.shapes.add_textbox(
                Emu(round(it["x"] * emu_px)), Emu(round(it["y"] * emu_px)),
                Emu(max(1, round(it["w"] * emu_px))), Emu(max(1, round(it["h"] * emu_px))))
            tf = tb.text_frame
            single = it["h"] <= 1.5 * it["fs"]
            tf.word_wrap = not single
            if hasattr(tf, 'auto_size'):
                tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP
            # Whitespace is collapsed only where the BROWSER collapses it.
            # `" ".join(text.split())` ran unconditionally, which is right under
            # `white-space: normal` (there a source newline renders as a space)
            # and lossy under `pre`, `pre-line` and `pre-wrap`, where the author
            # put the break in on purpose. The extractor never captured
            # `whiteSpace`, so the consumer could not tell the two apart; now it
            # does, and an authored break becomes a real paragraph.
            if str(it.get("ws", "normal")).startswith("pre"):
                lines = [ln.strip() for ln in it["text"].split("\n")]
                lines = [ln for ln in lines if ln] or [""]
            else:
                lines = [" ".join(it["text"].split())]
            for idx, line in enumerate(lines):
                para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                para.alignment = align_map.get(it["align"], PP_ALIGN.LEFT)
                if not single and it["lhr"]:
                    para.line_spacing = round(it["lhr"], 3)
                run = para.add_run()
                run.text = line
                run.font.size = Pt(it["fs"] * emu_px / 12700)
                run.font.name = it["fam"]
                run.font.italic = it["italic"]
                run.font.color.rgb = _rgb(it["color"])
    prs.save(str(pptx_path))
    n_fonts = embed_fonts(pptx_path, fam_italic, font_dirs, verbose=verbose) if font_dirs else 0
    return len(data), n_fonts


def _embeddable_font(family: str, font_dirs):
    """Best TTF/OTF file for a typeface name. PowerPoint embeds raw TTF/OTF only
    (never woff/woff2), so this resolver is restricted to those extensions,
    preferring TTF. Primary match is exact-token / oblique-intent (like
    _find_font_file), which handles separator-rich names ('GT-Standard-L-Standard-
    Medium'). Fallback is normalized-equality (_norm), which handles glued filenames
    whose tokens carry no separators ('31CHorizontalT03-560' for typeface
    '31C Horizontal T03 560')."""
    ttf_otf = [f for d in font_dirs if d.exists() for f in d.rglob("*")
               if f.suffix.lower() in (".ttf", ".otf")]
    fam_tokens = [t for t in re.split(r"[^a-z0-9]+", family.lower()) if t]
    fam_oblique = any(t in ("oblique", "italic") for t in fam_tokens)
    fam_core = {t for t in fam_tokens if t not in ("oblique", "italic")}

    best = None
    for f in ttf_otf:
        ft = {t for t in re.split(r"[^a-z0-9]+", f.stem.lower()) if t}
        if not fam_core.issubset(ft):
            continue
        if bool(ft & {"oblique", "italic"}) != fam_oblique:
            continue
        rank = (0 if f.suffix.lower() == ".ttf" else 1, len(ft - fam_core), len(f.stem))
        if best is None or rank < best[0]:
            best = (rank, f)
    if best:
        return best[1]

    # fallback: whole-name normalized equality (glued filenames)
    fam_norm = _norm(family)
    eq = [f for f in ttf_otf if _norm(f.stem) == fam_norm]
    eq.sort(key=lambda f: (0 if f.suffix.lower() == ".ttf" else 1, len(f.stem)))
    return eq[0] if eq else None


def embed_fonts(pptx_path: Path, fam_italic: dict, font_dirs, verbose=False):
    """Embed the used brand typefaces into the .pptx package so it renders
    identically on a machine without the fonts installed - the PowerPoint 'Embed
    fonts in the file' feature, applied at the OPC layer since python-pptx has no
    API for it. `fam_italic` maps each run typeface to whether any italic run uses
    it. Only TTF/OTF can be embedded. Returns the number of font files added."""
    import zipfile
    from lxml import etree

    # We only ever parse the OOXML parts python-pptx just wrote from a fixed template
    # (trusted, in-process - never untrusted external XML). A hardened parser with no
    # entity expansion and no network access makes the bandit B320 blacklist a
    # documented false positive rather than a real XXE/billion-laughs surface.
    xmlp = etree.XMLParser(resolve_entities=False, no_network=True)

    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    CT = "http://schemas.openxmlformats.org/package/2006/content-types"
    PR = "http://schemas.openxmlformats.org/package/2006/relationships"
    RT_FONT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"

    # Resolve one embeddable file per (typeface, style) slot.
    slots = []  # (typeface, 'regular'|'italic', Path)
    for fam in sorted(fam_italic):
        reg = _embeddable_font(fam, font_dirs)
        if reg:
            slots.append((fam, "regular", reg))
        elif verbose:
            print(f"{YELLOW}  embed: no TTF/OTF for typeface {fam}{RESET}")
        if fam_italic[fam]:
            ital = _embeddable_font(fam + " oblique", font_dirs) or _embeddable_font(fam + " italic", font_dirs)
            if ital:
                slots.append((fam, "italic", ital))
    if not slots:
        return 0

    with zipfile.ZipFile(pptx_path, "r") as z:
        members = {n: z.read(n) for n in z.namelist()}

    # 1. [Content_Types].xml: default mapping for the fntdata extension
    ctree = etree.fromstring(members["[Content_Types].xml"], xmlp)  # nosec B320 - trusted in-process OOXML
    if not any(d.get("Extension") == "fntdata" for d in ctree.findall(f"{{{CT}}}Default")):
        d = etree.SubElement(ctree, f"{{{CT}}}Default")
        d.set("Extension", "fntdata")
        d.set("ContentType", "application/x-fontdata")
    members["[Content_Types].xml"] = etree.tostring(
        ctree, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 2. font parts + presentation relationships
    rels_name = "ppt/_rels/presentation.xml.rels"
    rtree = etree.fromstring(members[rels_name], xmlp)  # nosec B320 - trusted in-process OOXML
    used_ids = {r.get("Id") for r in rtree.findall(f"{{{PR}}}Relationship")}

    def _new_id(seed):
        rid = f"rIdF{seed}"
        while rid in used_ids:
            seed += 1
            rid = f"rIdF{seed}"
        used_ids.add(rid)
        return rid

    slot_ids = []  # (typeface, style, rId)
    for i, (fam, style, path) in enumerate(slots, 1):
        members[f"ppt/fonts/font{i}.fntdata"] = path.read_bytes()
        rid = _new_id(i)
        rel = etree.SubElement(rtree, f"{{{PR}}}Relationship")
        rel.set("Id", rid)
        rel.set("Type", RT_FONT)
        rel.set("Target", f"fonts/font{i}.fntdata")
        slot_ids.append((fam, style, rid))
    members[rels_name] = etree.tostring(
        rtree, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 3. presentation.xml: embedTrueTypeFonts + embeddedFontLst (schema-ordered)
    ptree = etree.fromstring(members["ppt/presentation.xml"], xmlp)  # nosec B320 - trusted in-process OOXML
    ptree.set("embedTrueTypeFonts", "1")
    by_fam = {}
    for fam, style, rid in slot_ids:
        by_fam.setdefault(fam, {})[style] = rid
    lst = etree.Element(f"{{{P}}}embeddedFontLst", nsmap={"p": P, "r": R})
    for fam in sorted(by_fam):
        ef = etree.SubElement(lst, f"{{{P}}}embeddedFont")
        etree.SubElement(ef, f"{{{P}}}font").set("typeface", fam)
        for style in ("regular", "italic"):  # schema order: regular before italic
            if style in by_fam[fam]:
                etree.SubElement(ef, f"{{{P}}}{style}").set(f"{{{R}}}id", by_fam[fam][style])
    # CT_Presentation order: ... sldIdLst, sldSz, notesSz, smartTags, embeddedFontLst ...
    anchor = next((ptree.find(f"{{{P}}}{t}") for t in ("notesSz", "sldSz", "sldIdLst")
                   if ptree.find(f"{{{P}}}{t}") is not None), None)
    if anchor is not None:
        anchor.addnext(lst)
    else:
        ptree.append(lst)
    members["ppt/presentation.xml"] = etree.tostring(
        ptree, xml_declaration=True, encoding="UTF-8", standalone=True)

    # 4. rewrite the package, atomically.
    #
    # `ZipFile(pptx_path, "w")` TRUNCATES on open, and the path is the deck
    # `prs.save` wrote seconds earlier. Any failure inside the write loop -- a
    # full disk, an unreadable font file, an interrupt -- left a truncated
    # .pptx and no copy of the original, so the whole chromium render pass had
    # to run again. (The background PNGs survive under `editable-bg/`; the deck
    # does not.) `scripts/utils/atomic` is this repo's convention for exactly
    # this, and `partner-scorecard.py` already routes its rewrite through it.
    tmp = pptx_path.with_name(pptx_path.name + ".tmp")
    try:
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as z:
            for n, b in members.items():
                z.writestr(n, b)
        os.replace(tmp, pptx_path)
    except BaseException:
        tmp.unlink(missing_ok=True)
        raise
    if verbose:
        print(f"{GRAY}  embedded {len(slots)} font file(s): {', '.join(sorted(by_fam))}{RESET}")
    return len(slots)


def inline_html(work_html: Path, out_html: Path):
    """Portable single HTML: inline every remaining url()/src (images) as base64.
    Fonts are already inlined by resolve_fonts; this handles image fills/logos."""
    html = work_html.read_text(encoding="utf-8")
    base = work_html.parent

    def _uri(rel):
        fp = (base / rel).resolve()
        if not fp.exists():
            return None
        return _data_uri(fp)

    # Both patterns below used to accept one spelling each - `src="..."` only,
    # and `url("...")` / `url('...')` only. `src='images/a.png'` and the
    # unquoted `url(images/a.png)` are both valid and both survived untouched,
    # so the file `main` announces as "(self-contained)" 404'd its images the
    # moment it moved off the export directory. Measured 2026-08-30: both forms
    # still referenced `images/a.png` in the written output.
    def _src(m):
        quote = m.group(1)
        u = _uri(m.group(2))
        return f"src={quote}{u}{quote}" if u else m.group(0)

    html = re.sub(
        r"""src=(["'])(?!data:|https?:|file:)([^"']+)\1""", _src, html,
    )

    def _url(m):
        quote = m.group(1) or ""
        path = m.group("quoted") if m.group("quoted") is not None else m.group("bare")
        if path.startswith(("data:", "http", "file:")):
            return m.group(0)
        u = _uri(path)
        return f"url({quote}{u}{quote})" if u else m.group(0)

    html = re.sub(
        r"""url\(\s*(?:(["'])(?P<quoted>[^"']+)\1|(?P<bare>[^"')\s]+))\s*\)""",
        _url, html,
    )
    out_html.write_text(html, encoding="utf-8")
    return out_html


# ============================================================
# Main
# ============================================================

def main():
    ap = argparse.ArgumentParser(description="Render a Pencil export_html deck to PNG/PDF/PPTX/HTML.")
    ap.add_argument("--html", required=True, help="HTML produced by the Pencil export_html MCP tool")
    ap.add_argument("--out-dir", required=True, help="output directory (png/ + deliverables land here)")
    ap.add_argument("--fonts-dir", action="append", default=[], help="dir to search for brand fonts (repeatable)")
    ap.add_argument("--stem", default=None, help="output filename stem (default: HTML stem)")
    ap.add_argument("--formats", default="png,pdf,pptx,html",
                    help="comma list: png,pdf,html,pptx,pptx-flat. 'pptx' (default) emits the "
                         "EDITABLE <stem>.pptx (text-less backgrounds + native text boxes, brand "
                         "fonts embedded); 'pptx-flat' (alias pptx-image) emits the locked-look "
                         "image-per-slide '<stem> (ready to be shared with the world).pptx'. "
                         "'editable' is an alias of 'pptx'")
    ap.add_argument("--keep-in-bg", action="append", default=[],
                    help="editable mode: extra data-pencil-name whose text stays baked in the "
                         "background (branding/decoration), added to the built-in decor set (repeatable)")
    ap.add_argument("--width", type=int, default=1920)
    ap.add_argument("--height", type=int, default=1080)
    ap.add_argument("--scale", type=int, default=2)
    ap.add_argument("--image-format", choices=["png", "jpeg"], default="png",
                    help="per-slide render format. 'jpeg' compresses PDF/PPTX/HTML in one pass "
                         "(no NXPowerLite needed); 'png' is lossless (default)")
    ap.add_argument("--quality", type=int, default=85, help="JPEG quality 1-100 (only with --image-format jpeg)")
    ap.add_argument("--verbose", "-v", action="store_true")
    args = ap.parse_args()

    src_html = Path(args.html).resolve()
    if not src_html.exists():
        raise SystemExit(f"{RED}HTML not found: {src_html}{RESET}")
    out_dir = Path(args.out_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    png_dir = out_dir / "png"
    stem = args.stem or src_html.stem
    fonts = [Path(d).resolve() for d in args.fonts_dir]
    formats = {f.strip() for f in args.formats.split(",") if f.strip()}

    print(f"{BOLD}{CYAN}Pencil deck export{RESET}  {GRAY}{src_html.name} -> {out_dir}{RESET}")

    # Fonts: embed into a working copy placed NEXT TO the source so relative
    # image refs (images/..., ../assets/...) still resolve during chromium load.
    html = src_html.read_text(encoding="utf-8")
    if fonts:
        html = resolve_fonts(html, fonts, verbose=args.verbose)
    else:
        print(f"{YELLOW}  no --fonts-dir given; brand fonts may fall back{RESET}")
    work_html = src_html.parent / f"{stem}._render.html"
    work_html.write_text(html, encoding="utf-8")

    want_editable = bool(formats & {"pptx", "editable"})   # pptx == the editable twin
    want_flat = bool(formats & {"pptx-flat", "pptx-image"})  # opt-in locked-look image deck
    # render_pngs feeds png/pdf/html and the flat pptx (the editable twin renders its own bg)
    need_pngs = bool(formats & {"png", "pdf", "html"}) or want_flat
    try:
        if need_pngs:
            n = render_pngs(work_html, png_dir, args.width, args.height, args.scale,
                            img_format=args.image_format, quality=args.quality, verbose=args.verbose)
            print(f"{GREEN}  {args.image_format.upper():4}{RESET}  {n} slides -> {png_dir}")

            if "pdf" in formats:
                build_pdf(png_dir, out_dir / f"{stem}.pdf", args.width, args.height)
                print(f"{GREEN}  PDF{RESET}   {out_dir / (stem + '.pdf')}")
            if want_flat:
                flat_path = out_dir / f"{stem} (ready to be shared with the world).pptx"
                build_pptx(png_dir, flat_path)
                print(f"{GREEN}  FLAT{RESET}  {flat_path} (image-per-slide, portable, not editable)")
            if "html" in formats:
                inline_html(work_html, out_dir / f"{stem}.html")
                print(f"{GREEN}  HTML{RESET}  {out_dir / (stem + '.html')} (self-contained)")
            if "png" not in formats:
                for p in _slide_files(png_dir):
                    p.unlink()

        if want_editable:
            decor = list(DECOR_DEFAULT) + list(args.keep_in_bg)
            ebg = out_dir / "editable-bg"
            data = render_editable(work_html, ebg, args.width, args.height, args.scale,
                                   decor, img_format=args.image_format, quality=args.quality,
                                   verbose=args.verbose)
            edit_pptx = out_dir / f"{stem}.pptx"
            n_slides, n_fonts = build_editable_pptx(ebg, data, edit_pptx, args.width, args.height,
                                                    font_dirs=fonts, verbose=args.verbose)
            nb = sum(len(d["items"]) for d in data)
            fmsg = f", {n_fonts} font(s) embedded" if n_fonts else " (no fonts embedded)"
            print(f"{GREEN}  PPTX{RESET}  {edit_pptx} (editable: {n_slides} slides, {nb} text boxes{fmsg})")
    finally:
        work_html.unlink(missing_ok=True)

    print(f"{BOLD}{GREEN}done{RESET}")


if __name__ == "__main__":
    main()
