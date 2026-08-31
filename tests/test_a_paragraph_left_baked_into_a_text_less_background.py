#!/usr/bin/env python3
"""Shard 11-p1: five findings, and one background that was not text-less.

`pencil-export --format pptx` promises an EDITABLE deck: the Pencil render as a
text-less background image, with native PowerPoint text boxes on top. The
extractor took "leaf nodes only" -- `childElementCount > 0 -> skip` -- so a
paragraph carrying ANY inline markup was skipped. `<p>Our <strong>2026</strong>
plan</p>` has one element child, so the `<p>` was never extracted, never marked
`data-ov`, and never hidden: the words around the bold stayed painted into the
image while `<strong>` alone became a floating text box. Editing the deck could
not reach them, and moving the box left "Our" and "plan" behind.

The tool then printed "N text-less backgrounds", which is a claim about the whole
frame that this pass never established -- and that sentence is what made the
partial extraction invisible. The measurement is now real (`leftover`) and the
line says what it found.

Two more in the same file: whitespace was collapsed unconditionally, which is
correct under `white-space: normal` (the browser renders a source newline as a
space) and lossy under `pre`/`pre-line`/`pre-wrap`, where the author put the
break in deliberately -- and the extractor never captured `whiteSpace`, so the
consumer could not tell them apart. And `embed_fonts` reopened the just-saved
.pptx with mode `"w"`, which TRUNCATES on open, so any failure in the write loop
left a corrupt package and no copy of the deck.

Then two elsewhere. `output-organizer archive` filtered on `f.parts` of an
ABSOLUTE path, so one directory called `archive` anywhere in the workspace's
ancestry excluded every file -- and it reported that in green as "No files older
than N days found." And `partner-scorecard` did not recognise a table separator
carrying alignment colons (`|:---|---:|`), so every data row was skipped.

Run: .venv/bin/python -m pytest tests/test_a_paragraph_left_baked_into_a_text_less_background.py
"""
from __future__ import annotations

import importlib.util
import struct
import sys
import zipfile
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))


def _load(rel: str, name: str):
    spec = importlib.util.spec_from_file_location(name, ROOT / rel)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


px = _load("scripts/pencil-export.py", "pencil_export_11p1")
oo = _load("scripts/output-organizer.py", "output_organizer_11p1")
ps = _load("scripts/partner-scorecard.py", "partner_scorecard_11p1")


# ============================================================
# Finding 2 -- the paragraph the extractor skipped
# ============================================================

SLIDE = """
<html><body style="margin:0">
<div data-pencil-name="Slide-1" data-pencil-id="s1"
     style="position:relative;width:600px;height:300px;font-family:Arial">
  {body}
</div>
</body></html>
"""


@pytest.fixture(scope="module")
def browser():
    sync_playwright = pytest.importorskip("playwright.sync_api").sync_playwright
    with sync_playwright() as p:
        b = p.chromium.launch()
        yield b
        b.close()


def extract(browser, body: str, decor=()):
    """Run the REAL EXTRACT_JS in a real browser over one slide."""
    page = browser.new_page()
    try:
        page.set_content(SLIDE.format(body=body))
        return page.evaluate(px.EXTRACT_JS, list(decor))[0]
    finally:
        page.close()


@pytest.mark.requires_playwright
def test_a_paragraph_with_inline_markup_is_extracted_whole(browser):
    """The regression. `<p>` had one element child, so the leaf rule skipped it
    and only `<strong>` came out -- the words around the bold stayed in the
    image."""
    slide = extract(browser, '<p style="font-size:14px">Our <strong>2026</strong> plan</p>')
    texts = [i["text"].strip() for i in slide["items"]]
    assert texts == ["Our 2026 plan"]


@pytest.mark.requires_playwright
def test_the_inline_child_is_not_extracted_a_second_time(browser):
    """Hiding the parent hides the child, so extracting both would put the same
    words on the slide twice."""
    slide = extract(browser, '<p style="font-size:14px">Our <strong>2026</strong> plan</p>')
    assert len(slide["items"]) == 1


@pytest.mark.requires_playwright
def test_an_inline_block_child_does_not_split_the_paragraph(browser):
    """`inline-block` is how a deck writes a badge or a pill inside a line.
    Treating it as a block would push the paragraph back into the skipped
    branch and bake its words into the image again."""
    slide = extract(
        browser,
        '<p style="font-size:14px">Stage '
        '<span style="display:inline-block;padding:2px">2</span> of 5</p>')
    assert [i["text"].strip() for i in slide["items"]] == ["Stage 2 of 5"]
    assert slide["leftover"] == 0


@pytest.mark.requires_playwright
def test_two_block_paragraphs_stay_two_boxes(browser):
    """The container must NOT be taken as one paragraph: they are laid out
    separately and their boxes are what makes them editable separately."""
    slide = extract(browser, '<div><p style="font-size:14px">First</p>'
                             '<p style="font-size:14px">Second</p></div>')
    assert sorted(i["text"].strip() for i in slide["items"]) == ["First", "Second"]


@pytest.mark.requires_playwright
def test_a_plain_leaf_still_works(browser):
    slide = extract(browser, '<span style="font-size:14px">Just words</span>')
    assert [i["text"].strip() for i in slide["items"]] == ["Just words"]


@pytest.mark.requires_playwright
def test_decorative_text_still_stays_in_the_background(browser):
    """`--keep-in-bg` names must not be dragged out by the wider rule."""
    slide = extract(
        browser,
        '<div data-pencil-name="Footer"><p style="font-size:10px">31 Concept</p></div>'
        '<p style="font-size:14px">Content</p>',
        decor=("Footer",))
    assert [i["text"].strip() for i in slide["items"]] == ["Content"]
    assert slide["leftover"] == 0, "decorative text is excluded from the count"


@pytest.mark.requires_playwright
def test_text_that_could_not_be_extracted_is_COUNTED(browser):
    """An icon-font glyph is deliberately left in the background. The point is
    that the tool now KNOWS, instead of calling the result text-less."""
    slide = extract(
        browser,
        '<p style="font-size:14px;font-family:Material Icons">home</p>'
        '<p style="font-size:14px">Content</p>')
    assert [i["text"].strip() for i in slide["items"]] == ["Content"]
    assert slide["leftover"] == 1


@pytest.mark.requires_playwright
def test_a_fully_extracted_slide_reports_nothing_left(browser):
    slide = extract(browser, '<p style="font-size:14px">Our <em>whole</em> point</p>')
    assert slide["leftover"] == 0


@pytest.mark.requires_playwright
def test_the_white_space_mode_is_captured(browser):
    """Finding 3's root cause: the consumer had no way to tell an authored break
    from source indentation, because this was never recorded."""
    slide = extract(browser, '<p style="font-size:14px;white-space:pre">a\nb</p>')
    assert slide["items"][0]["ws"].startswith("pre")


# ============================================================
# Finding 2b -- the sentence that covered for the gap
# ============================================================

def test_a_complete_extraction_no_longer_claims_a_text_less_background():
    msg, incomplete = px.extraction_summary(
        [{"items": [1, 2], "leftover": 0}, {"items": [3], "leftover": 0}])
    assert incomplete is False
    assert "text-less" not in msg
    assert "no content text left in them" in msg


def test_an_incomplete_extraction_says_so_and_counts_it():
    msg, incomplete = px.extraction_summary(
        [{"items": [1], "leftover": 2}, {"items": [], "leftover": 1}])
    assert incomplete is True
    assert "3 text node(s) could NOT be extracted" in msg


def test_an_older_payload_without_the_field_does_not_crash_the_summary():
    """Defensive on purpose: `data` also arrives from a cached run."""
    msg, incomplete = px.extraction_summary([{"items": [1]}])
    assert incomplete is False


# ============================================================
# Finding 3 -- whitespace collapsed where the browser would not
# ============================================================

def _png(path: Path):
    """A 1x1 PNG, so `add_picture` has something real to embed."""
    def chunk(tag, data):
        c = tag + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c))
    import zlib
    raw = (b"\x89PNG\r\n\x1a\n"
           + chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
           + chunk(b"IDAT", zlib.compress(b"\x00\x00\x00\x00"))
           + chunk(b"IEND", b""))
    path.write_bytes(raw)
    return path


def _item(text, ws="normal", **over):
    base = {"text": text, "x": 0, "y": 0, "w": 300, "h": 100, "fs": 14,
            "fam": "Arial", "italic": False, "color": "rgb(0, 0, 0)",
            "align": "left", "lhr": 1.2, "ws": ws}
    base.update(over)
    return base


_BUILDS = [0]


def _build(tmp_path, items):
    pytest.importorskip("pptx")
    _BUILDS[0] += 1
    bg = tmp_path / f"bg{_BUILDS[0]}"
    bg.mkdir()
    _png(bg / "slide-01.png")
    out = tmp_path / f"deck{_BUILDS[0]}.pptx"
    px.build_editable_pptx(bg, [{"name": "Slide-1", "id": "s1", "items": items}],
                           out, 1920, 1080)
    from pptx import Presentation
    slide = Presentation(str(out)).slides[0]
    boxes = [sh for sh in slide.shapes if sh.has_text_frame]
    return boxes


def test_an_authored_break_survives_as_a_real_paragraph(tmp_path):
    boxes = _build(tmp_path, [_item("first line\nsecond line", ws="pre")])
    paras = [p.text for p in boxes[0].text_frame.paragraphs]
    assert paras == ["first line", "second line"]


def test_pre_wrap_and_pre_line_are_treated_the_same_way(tmp_path):
    for mode in ("pre-line", "pre-wrap"):
        boxes = _build(tmp_path, [_item("a\nb", ws=mode)])
        assert [p.text for p in boxes[0].text_frame.paragraphs] == ["a", "b"]


def test_source_indentation_is_still_collapsed(tmp_path):
    """Under `white-space: normal` the browser renders those newlines as spaces,
    so collapsing them is CORRECT, not lossy. The fix must not break it."""
    boxes = _build(tmp_path, [_item("  Our\n   long\n  title  ")])
    paras = [p.text for p in boxes[0].text_frame.paragraphs]
    assert paras == ["Our long title"]


def test_a_missing_white_space_field_falls_back_to_collapsing(tmp_path):
    item = _item("a\nb")
    del item["ws"]
    boxes = _build(tmp_path, [item])
    assert [p.text for p in boxes[0].text_frame.paragraphs] == ["a b"]


def test_an_empty_pre_block_still_gets_a_styled_paragraph(tmp_path):
    """`or [""]` is not decoration. With no lines the loop never runs, so
    paragraph 0 keeps python-pptx's default alignment instead of the extracted
    one -- the same box would render differently from every non-pre sibling."""
    from pptx.enum.text import PP_ALIGN
    boxes = _build(tmp_path, [_item("\n\n", ws="pre", align="center")])
    para = boxes[0].text_frame.paragraphs[0]
    assert para.text == ""
    assert para.alignment == PP_ALIGN.CENTER


# ============================================================
# Finding 5 -- a rewrite that truncated first and wrote after
# ============================================================

def _pptx_with(tmp_path, name="deck.pptx"):
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    prs.slides.add_slide(prs.slide_layouts[6])
    p = tmp_path / name
    prs.save(str(p))
    return p


def test_a_failed_font_embed_leaves_the_deck_intact(tmp_path, monkeypatch):
    """`ZipFile(path, "w")` truncates on open, and the path is the deck saved
    seconds earlier. A failure in the write loop used to leave a corrupt package
    and no original, costing a full chromium re-render."""
    deck = _pptx_with(tmp_path)
    before = deck.read_bytes()

    real_writestr = zipfile.ZipFile.writestr
    state = {"n": 0}

    def flaky(self, *a, **k):
        state["n"] += 1
        if state["n"] == 3:
            raise OSError(28, "No space left on device")
        return real_writestr(self, *a, **k)
    monkeypatch.setattr(zipfile.ZipFile, "writestr", flaky)

    font = tmp_path / "fonts" / "Arial.ttf"
    font.parent.mkdir()
    font.write_bytes(b"\x00\x01\x00\x00not-a-real-font")

    with pytest.raises(OSError):
        px.embed_fonts(deck, {"Arial": False}, [font.parent])

    assert deck.read_bytes() == before, "the deck was truncated in place"


def test_no_temporary_file_is_left_behind_after_a_failure(tmp_path, monkeypatch):
    deck = _pptx_with(tmp_path)
    real_writestr = zipfile.ZipFile.writestr
    state = {"n": 0}

    def flaky(self, *a, **k):
        state["n"] += 1
        if state["n"] == 3:
            raise OSError(28, "No space left on device")
        return real_writestr(self, *a, **k)
    monkeypatch.setattr(zipfile.ZipFile, "writestr", flaky)

    font = tmp_path / "fonts" / "Arial.ttf"
    font.parent.mkdir()
    font.write_bytes(b"\x00\x01\x00\x00not-a-real-font")
    with pytest.raises(OSError):
        px.embed_fonts(deck, {"Arial": False}, [font.parent])
    assert list(tmp_path.glob("*.tmp")) == []


def test_a_successful_embed_still_replaces_the_deck(tmp_path):
    deck = _pptx_with(tmp_path)
    font = tmp_path / "fonts" / "Arial.ttf"
    font.parent.mkdir()
    font.write_bytes(b"\x00\x01\x00\x00not-a-real-font")
    assert px.embed_fonts(deck, {"Arial": False}, [font.parent]) == 1
    with zipfile.ZipFile(deck) as z:
        assert any(n.startswith("ppt/fonts/") for n in z.namelist())


# ============================================================
# Finding 1 -- an archive that reported success for doing nothing
# ============================================================

@pytest.fixture
def outputs(tmp_path, monkeypatch):
    """An outputs tree under an ancestor directory literally named `archive`."""
    root = tmp_path / "archive" / "workspace" / "outputs"
    root.mkdir(parents=True)
    monkeypatch.setattr(oo, "outputs_dir", lambda p=root: p)
    return root


def _age(path: Path, days: int):
    import os
    import time
    t = time.time() - days * 86400
    os.utime(path, (t, t))


def test_an_ancestor_named_archive_no_longer_hides_every_file(outputs, capsys):
    """`f.parts` is the WHOLE absolute path, so one `archive` anywhere above the
    workspace excluded the entire tree -- and the tool said so in green."""
    old = outputs / "documents" / "old.md"
    old.parent.mkdir(parents=True)
    old.write_text("x", encoding="utf-8")
    _age(old, 90)

    oo.archive(30, execute=False)
    out = capsys.readouterr().out
    assert "No files older than" not in out
    assert "old.md" in out


def test_the_real_archive_subtree_is_still_excluded(outputs, capsys):
    """The filter's actual job: files already archived must not be re-archived."""
    done = outputs / "archive" / "done.md"
    done.parent.mkdir(parents=True)
    done.write_text("x", encoding="utf-8")
    _age(done, 90)

    oo.archive(30, execute=False)
    assert "No files older than" in capsys.readouterr().out


def test_a_recent_file_is_still_not_archived(outputs, capsys):
    fresh = outputs / "fresh.md"
    fresh.write_text("x", encoding="utf-8")
    oo.archive(30, execute=False)
    assert "No files older than" in capsys.readouterr().out


def test_the_dry_run_moves_nothing(outputs, capsys):
    old = outputs / "old.md"
    old.write_text("x", encoding="utf-8")
    _age(old, 90)
    oo.archive(30, execute=False)
    assert old.is_file()


# ============================================================
# Finding 4 -- a separator row with alignment colons
# ============================================================

def _table(sep: str) -> str:
    return (
        "## Partnership Discussions\n\n"
        "| Partner | Stage | Value | Next Step | Owner | Notes | Health |\n"
        f"{sep}\n"
        "| Acme | Pilot | 100 | Call | misha | none | green |\n"
    )


def test_a_plain_separator_still_parses():
    rows = ps.parse_partnerships(_table("|---|---|---|---|---|---|---|"))
    assert [r["partner"] for r in rows] == ["Acme"]


@pytest.mark.parametrize("sep", [
    "|:---|:---|:---|:---|:---|:---|:---|",
    "|---:|---:|---:|---:|---:|---:|---:|",
    "|:---:|:---:|:---:|:---:|:---:|:---:|:---:|",
    "|:-------|--------:|:------:|--------|--------|--------|--------|",
])
def test_an_aligned_separator_parses_too(sep):
    """GitHub-flavoured alignment colons. `startswith("---")` matched none of
    them, so every data row was skipped and the tool exited 2 over a table that
    was right there."""
    rows = ps.parse_partnerships(_table(sep))
    assert [r["partner"] for r in rows] == ["Acme"]


def test_a_row_that_only_looks_like_a_separator_is_not_one():
    """The pattern is anchored: a real value starting with a dash must not be
    mistaken for the separator and swallow the row."""
    body = _table("|---|---|---|---|---|---|---|").replace(
        "| Acme |", "| ---not-a-sep |")
    rows = ps.parse_partnerships(body)
    assert [r["partner"] for r in rows] == ["---not-a-sep"]


def test_a_dashed_VALUE_before_the_separator_is_not_taken_as_one():
    """Anchoring matters before the real separator row, not after it. With a
    prefix match, `| ---weird |` flips `seen_sep` early, and the genuine
    separator that follows is then parsed as a data row."""
    body = (
        "## Partnership Discussions\n\n"
        "| Partner | Stage | Value | Next | Owner | Notes | Health |\n"
        "| ---weird | a | b | c | d | e | f |\n"
        "|---|---|---|---|---|---|---|\n"
        "| Acme | Pilot | 100 | Call | misha | none | green |\n"
    )
    rows = ps.parse_partnerships(body)
    assert [r["partner"] for r in rows] == ["Acme"], (
        "the dashed value was accepted as the separator")


def test_a_two_dash_separator_is_not_accepted():
    """`{3,}` is the GFM minimum. Loosening it would let a stray `--` cell act
    as a separator."""
    rows = ps.parse_partnerships(_table("|--|--|"))
    assert rows == []
