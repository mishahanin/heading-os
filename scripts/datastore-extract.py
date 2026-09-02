#!/usr/bin/env python3
"""
DataStore Extraction Script for 31C Workspace

Converts binary files (PPTX, XLSX, PDF, DOCX) in the datastore/ to readable
markdown companion files (-extract.md). The companion is the only surface the
memory index can reach, so a document type this script does not handle is a
document the workspace cannot recall at all.

Usage:
    python scripts/datastore-extract.py                        # scan and extract all new files
    python scripts/datastore-extract.py datastore/investment/ceo-only/    # extract from specific folder
    python scripts/datastore-extract.py --update-index          # also update INDEX.md
    python scripts/datastore-extract.py --force                 # re-extract even if companion exists

Prerequisites:
    openpyxl, python-pptx, python-docx, pdfminer.six

    Measured 2026-09-02: the first three are pinned in `pyproject.toml`.
    pdfminer.six is NOT declared anywhere in this project; it is installed only
    as a transitive dependency of `markitdown[pdf]`. Dropping or narrowing that
    extra silently removes PDF extraction, and no pin in this repository will
    fail to warn you.

Tests:
    tests/test_a_rollback_that_deleted_what_it_never_backed_up.py
    tests/test_an_extractor_blind_to_two_thirds_of_the_corpus.py
"""

import argparse
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import NamedTuple

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from scripts.utils.colors import GREEN, YELLOW, RED, BOLD, RESET
from scripts.utils.workspace import get_workspace_root, get_datastore_dir, get_default_tz

WORKSPACE = get_workspace_root()


def datastore_dir():
    """Resolved at call time, never at import.

    `get_datastore_dir()` reads `HEADING_OS_DATA` on every call, so it follows
    the environment for a caller that asks after the environment moved. As a
    module-level constant it asked once, during its own import, and stored the
    answer, so a test that imported this module and then repointed the data
    root still read and wrote the operator's real overlay. `index_file()` below
    derived from it, so freezing the first froze the second too.
    """
    return get_datastore_dir()


def index_file():
    return datastore_dir() / "INDEX.md"


def extract_xlsx(filepath):
    """Extract XLSX content to markdown."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        print(f"{RED}Error: openpyxl not installed. Run: pip install openpyxl{RESET}")
        return None

    wb = load_workbook(filepath, data_only=True)
    lines = []
    lines.append(f"# Extract: {filepath.name}")
    lines.append(f"")
    lines.append(f"> Auto-extracted from `{filepath.name}` on {datetime.now(get_default_tz()).strftime('%Y-%m-%d')}")
    lines.append(f"> This is a companion file for Claude to read. The original XLSX is the source of truth.")
    lines.append(f"")

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## Sheet: {sheet_name}")
        lines.append(f"")

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            lines.append("*(empty sheet)*")
            lines.append("")
            continue

        # Find header row (first non-empty row)
        header_row = None
        data_start = 0
        for i, row in enumerate(rows):
            if any(cell is not None for cell in row):
                header_row = row
                data_start = i + 1
                break

        if header_row is None:
            lines.append("*(empty sheet)*")
            lines.append("")
            continue

        # Build markdown table
        headers = [_cell(h) for h in header_row]
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

        for row in rows[data_start:]:
            cells = [_cell(c) for c in row]
            # Truncate long cells
            cells = [c[:100] + "..." if len(c) > 100 else c for c in cells]
            # Pad to header length
            while len(cells) < len(headers):
                cells.append("")
            lines.append("| " + " | ".join(cells[:len(headers)]) + " |")

        lines.append("")

    return "\n".join(lines)


def _cell(value) -> str:
    """One spreadsheet cell, safe inside a markdown table row.

    Cells were joined straight into `" | "`, so any value containing a pipe
    added a phantom column and any value containing a newline ended the row
    early — and the companion `-extract.md` stopped representing the sheet it
    claims to be an extract of.
    """
    if value is None:
        return ""
    text = str(value)
    return text.replace("\\", "\\\\").replace("|", "\\|").replace("\r", " ").replace("\n", " ")


def extract_pptx(filepath):
    """Extract PPTX text content to markdown."""
    try:
        from pptx import Presentation
    except ImportError:
        print(f"{RED}Error: python-pptx not installed. Run: pip install python-pptx{RESET}")
        return None

    prs = Presentation(str(filepath))
    lines = []
    lines.append(f"# Extract: {filepath.name}")
    lines.append(f"")
    lines.append(f"> Auto-extracted from `{filepath.name}` on {datetime.now(get_default_tz()).strftime('%Y-%m-%d')}")
    lines.append(f"> This is a companion file for Claude to read. The original PPTX is the source of truth.")
    lines.append(f"")
    lines.append(f"Total slides: {len(prs.slides)}")
    lines.append(f"")

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        lines.append(f"")

        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)

        if texts:
            for text in texts:
                lines.append(f"- {text}")
        else:
            lines.append("*(no text content - may be an image/diagram slide)*")

        lines.append("")

    return "\n".join(lines)


# The ceiling on one companion's extracted text, in characters.
#
# 100,000 characters is roughly 15,000 words, or about 25,000 tokens: larger
# than any hand-written note in this workspace, and still small enough that an
# agent can read the whole companion in one go. The number exists because the
# companion feeds a shared memory index. A single 900-page manual extracted
# whole would out-mass every other document in the datastore put together, and
# every unrelated query would start retrieving chunks of it. Truncation is the
# lesser harm as long as it is VISIBLE, which is why the marker is written at
# both ends of the companion rather than only where the cut happened.
MAX_EXTRACT_CHARS = 100_000


def _extract_header(filepath, kind: str) -> list:
    """The header every companion carries, whatever the source format was."""
    return [
        f"# Extract: {filepath.name}",
        "",
        f"> Auto-extracted from `{filepath.name}` on "
        f"{datetime.now(get_default_tz()).strftime('%Y-%m-%d')}",
        f"> This is a companion file for Claude to read. The original {kind} "
        f"is the source of truth.",
        "",
    ]


def _no_text_companion(filepath, kind: str, detail: str) -> str:
    """A companion for a binary that yielded no text at all.

    The alternative (write nothing) was rejected. A companion holding only the
    header is worse than useless: it reads to the next agent as "this document
    is empty", which is a false statement about a file that may be the most
    important one in its folder. Writing no file at all is quieter but leaves
    the same document unreachable forever, with every later run re-opening the
    same scan to re-discover the same nothing.

    So the companion is written and SAYS SO. It is indexable, it names the
    binary as the thing to open, and `--force` re-extracts it the day an OCR
    path exists.
    """
    lines = _extract_header(filepath, kind)
    lines.append(f"**No extractable text layer.** {detail}")
    lines.append("")
    lines.append(
        f"This {kind} carries no text this extractor can read. Open the "
        f"original `{filepath.name}` directly; it is the source of truth. This "
        f"companion is NOT the record of an empty document."
    )
    lines.append("")
    return "\n".join(lines)


def _truncation_note(cap: int, where: str, kind: str) -> str:
    return (f"> TRUNCATED at the {cap}-character cap, after {where}. This "
            f"companion is PARTIAL; open the original {kind} for the rest.")


def extract_pdf(filepath):
    """Extract PDF text to markdown, one section per page.

    A PDF with no text layer gets a companion that says so rather than an empty
    one or no file at all; the reasoning is in `_no_text_companion`.

    The extract is capped at `MAX_EXTRACT_CHARS` and cut at a PAGE boundary,
    never mid-page, so no section of the companion is a half sentence. The
    cap's own reasoning sits on the constant.

    Nothing is caught here. An encrypted PDF raises `PDFPasswordIncorrect` and
    a corrupt one raises `PDFSyntaxError`, both out of `extract_pages`, and
    `scan_and_extract` already names and counts a per-file parser failure
    without letting it end the batch. Catching them here would report the file
    as extracted-but-empty, which is the one outcome that must not happen.
    """
    try:
        from pdfminer.high_level import extract_pages
        from pdfminer.layout import LTTextContainer
    except ImportError:
        print(f"{RED}Error: pdfminer.six not installed. "
              f"Run: pip install pdfminer.six{RESET}")
        return None

    body: list = []
    total_chars = 0
    pages_read = 0
    truncated_after = None

    for page_number, page in enumerate(extract_pages(str(filepath)), 1):
        if total_chars >= MAX_EXTRACT_CHARS:
            truncated_after = page_number - 1
            break

        texts = []
        for element in page:
            if isinstance(element, LTTextContainer):
                text = element.get_text().strip()
                if text:
                    texts.append(text)

        pages_read = page_number
        if not texts:
            # An image-only page inside a document that has text elsewhere.
            continue

        body.append(f"## Page {page_number}")
        body.append("")
        for text in texts:
            body.append(text)
            body.append("")
        total_chars += sum(len(t) for t in texts)

    if not body:
        return _no_text_companion(
            filepath, "PDF",
            f"{pages_read} page(s) were read and none of them carried text, "
            f"which usually means a scan or an image-only document.")

    note = None
    if truncated_after is not None:
        note = _truncation_note(MAX_EXTRACT_CHARS,
                                f"page {truncated_after}", "PDF")

    lines = _extract_header(filepath, "PDF")
    if note:
        lines.append(note)
        lines.append("")
    lines.append(f"Pages read: {pages_read}")
    lines.append("")
    lines.extend(body)
    if note:
        lines.append(note)
        lines.append("")

    return "\n".join(lines)


def _docx_table_rows(table) -> list:
    """One DOCX table as markdown rows, the first row read as the header."""
    rows: list = []
    for row in table.rows:
        cells = [_cell(c.text) for c in row.cells]
        if not cells:
            continue
        cells = [c[:100] + "..." if len(c) > 100 else c for c in cells]
        rows.append("| " + " | ".join(cells) + " |")
        if len(rows) == 1:
            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return rows


def extract_docx(filepath):
    """Extract DOCX paragraphs AND tables to markdown, in document order.

    `Document.paragraphs` and `Document.tables` are separate collections. A
    reader that walks only the first returns nothing at all for a document
    whose content lives in tables (a rate card, a compliance matrix, a
    requirements sheet), and one that walks the second after the first reorders
    the document into two halves that were never adjacent. So the body's own
    child elements are walked instead, and a table is extracted where it
    actually sits.

    Capped at `MAX_EXTRACT_CHARS` and cut at a block boundary (a whole
    paragraph, or a whole table), never inside one.

    Nothing is caught here, for the reason given in `extract_pdf`.
    """
    try:
        import docx
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError:
        print(f"{RED}Error: python-docx not installed. "
              f"Run: pip install python-docx{RESET}")
        return None

    document = docx.Document(str(filepath))
    children = list(document.element.body.iterchildren())

    body: list = []
    total_chars = 0
    paragraphs = 0
    tables = 0
    truncated_after = None

    for index, child in enumerate(children, 1):
        if total_chars >= MAX_EXTRACT_CHARS:
            truncated_after = index - 1
            break

        if child.tag == qn("w:p"):
            text = Paragraph(child, document).text.strip()
            if not text:
                continue
            body.append(text)
            body.append("")
            total_chars += len(text)
            paragraphs += 1
        elif child.tag == qn("w:tbl"):
            rows = _docx_table_rows(Table(child, document))
            if not rows:
                continue
            body.extend(rows)
            body.append("")
            total_chars += sum(len(r) for r in rows)
            tables += 1

    if not body:
        return _no_text_companion(
            filepath, "DOCX",
            "The document body holds no paragraph text and no tables.")

    note = None
    if truncated_after is not None:
        note = _truncation_note(
            MAX_EXTRACT_CHARS,
            f"block {truncated_after} of {len(children)}", "DOCX")

    lines = _extract_header(filepath, "DOCX")
    if note:
        lines.append(note)
        lines.append("")
    lines.append(f"Paragraphs: {paragraphs}")
    lines.append(f"Tables: {tables}")
    lines.append("")
    lines.extend(body)
    if note:
        lines.append(note)
        lines.append("")

    return "\n".join(lines)


EXTRACTABLE_SUFFIXES = (".xlsx", ".pptx", ".pdf", ".docx")


def get_companion_path(filepath):
    """Get the -extract.md companion path for a binary file."""
    return filepath.with_name(filepath.stem + "-extract.md")


def _ambiguous_stems(paths) -> set:
    """(parent, stem) pairs claimed by more than one extractable file."""
    by_stem: dict = {}
    for path in paths:
        by_stem.setdefault((path.parent, path.stem), set()).add(path.suffix.lower())
    return {key for key, suffixes in by_stem.items() if len(suffixes) > 1}


def _companion_for(filepath, ambiguous: set):
    """The companion path, with the source suffix added only when it must be.

    `<stem>-extract.md` is derived from the STEM, so `pitch.pptx` and
    `pitch.xlsx` in one folder claim the same companion. The deck was extracted
    first (it sorts first), and then the workbook was told "Skip pitch.xlsx
    (companion already exists)" - a companion that describes the deck. The
    workbook was never extracted and nothing said so. Under `--force` it was
    worse: both ran, the second overwrote the first, and the run reported two
    files extracted with one file of output.

    Only a colliding pair is renamed. Every companion on disk today belongs to a
    file with no same-stem sibling, so none of them moves.
    """
    base = get_companion_path(filepath)
    if (filepath.parent, filepath.stem) not in ambiguous:
        return base
    suffix = filepath.suffix.lstrip(".").lower()
    return base.with_name(f"{filepath.stem}-{suffix}-extract.md")


def orphaned_companions(ambiguous: set) -> list:
    """Unsuffixed companions left behind when a same-stem sibling ARRIVES.

    The rename above is scoped to currently-ambiguous stems on the stated
    assumption that "every companion on disk today belongs to a file with no
    same-stem sibling". That is true on the day it was written and stops being
    true the moment somebody drops `pitch.xlsx` beside an already-extracted
    `pitch.pptx`: both sources now resolve to suffixed companions, and the old
    `pitch-extract.md` -- which describes only the deck -- sits between them,
    unreferenced, for the next reader to take as the companion for either. The
    stem was printed in the ambiguity warning; the stale FILE never was.

    Reported, never deleted. It is extracted content, the operator may have
    edited it, and which of the two files it describes is a question this
    script can answer only by guessing.
    """
    orphans = []
    for parent, stem in sorted(ambiguous):
        stale = parent / f"{stem}-extract.md"
        if stale.is_file():
            orphans.append(stale)
    return orphans


class ScanResult(NamedTuple):
    """What one scan produced, and what it could not."""
    extracted: list
    failures: list


def scan_and_extract(target_dir=None, force=False):
    """Scan for binary files and create companion extracts.

    Returns the extracted `(original, companion)` pairs and nothing else,
    because that is the contract `update_index` and three test modules already
    consume. A caller that must also know what FAILED calls `scan_report`
    instead. `main` is one, so that the process can exit non-zero rather than
    announce a run of 40 that produced 38 in the same words as one that
    produced 40.
    """
    return scan_report(target_dir, force).extracted


def scan_report(target_dir=None, force=False) -> ScanResult:
    """Scan for binary files and create companion extracts."""
    scan_dir = Path(target_dir) if target_dir else datastore_dir()

    if not scan_dir.exists():
        print(f"{RED}Directory not found: {scan_dir}{RESET}")
        return ScanResult([], [])

    # Find every extractable file, case-insensitively. `rglob("*.xlsx")` is
    # case-SENSITIVE on Linux, so `Q3.XLSX` was neither extracted nor
    # mentioned and the run printed its summary as though the datastore were
    # fully processed. On macOS and Windows the same pattern happened to match,
    # so the tool behaved differently depending on the filesystem underneath
    # it. One walk, one suffix test on the lowercased suffix.
    binary_files = sorted(
        p for p in scan_dir.rglob("*")
        if p.is_file() and p.suffix.lower() in EXTRACTABLE_SUFFIXES
    )

    if not binary_files:
        print(f"{YELLOW}No XLSX or PPTX files, and no PDF or DOCX files, "
              f"found in {scan_dir}{RESET}")
        return ScanResult([], [])

    ambiguous = _ambiguous_stems(binary_files)
    if ambiguous:
        # Named, because a renamed companion is a surprise the operator should
        # hear about once rather than discover later.
        print(f"{YELLOW}{len(ambiguous)} stem(s) are claimed by more than one "
              f"extractable file; those companions carry the source "
              f"suffix:{RESET}")
        for parent, stem in sorted(ambiguous):
            print(f"  {YELLOW}{parent.name}/{stem}{RESET}")
        orphans = orphaned_companions(ambiguous)
        if orphans:
            print(f"{YELLOW}{len(orphans)} unsuffixed companion(s) now describe "
                  f"only one of a pair. Read and delete or rename by hand:{RESET}")
            for stale in orphans:
                print(f"  {YELLOW}{stale}{RESET}")

    extracted = []
    failures: list[tuple[Path, str]] = []
    for filepath in sorted(binary_files):
        companion = _companion_for(filepath, ambiguous)

        if companion.exists() and not force:
            print(f"  {GREEN}Skip{RESET}  {filepath.name} (companion already exists)")
            continue

        print(f"  {BOLD}Extracting{RESET}  {filepath.name}...")

        # Per file, not per batch. `load_workbook`/`Presentation` raise on a
        # corrupt zip, an encrypted file or a parser error, and with nothing
        # catching them here the exception left the LOOP: one bad binary in a
        # datastore folder meant every later extractable file was never
        # processed, and the run ended on a traceback rather than a report.
        try:
            if filepath.suffix.lower() == ".xlsx":
                content = extract_xlsx(filepath)
            elif filepath.suffix.lower() == ".pptx":
                content = extract_pptx(filepath)
            elif filepath.suffix.lower() == ".pdf":
                content = extract_pdf(filepath)
            elif filepath.suffix.lower() == ".docx":
                content = extract_docx(filepath)
            else:
                continue
        except Exception as exc:  # noqa: BLE001 - one file's parser, not the batch
            print(f"  {RED}Failed{RESET}  {filepath.name}: "
                  f"{type(exc).__name__}: {exc}")
            failures.append((filepath, f"{type(exc).__name__}: {exc}"))
            continue

        if content:
            # The WRITE is per file too, and it was not. The try above covers
            # only the parsers, so an `OSError` out of `write_text` left the
            # loop exactly the way a parser exception used to: every later
            # extractable file went unprocessed, the run ended on a traceback,
            # and the `failures` report below was never printed. A read-only
            # datastore folder, a full disk, or a path-length limit on the
            # renamed `{stem}-{suffix}-extract.md` form is at least as common an
            # operator condition as a corrupt zip, and the tool recovered from
            # the corrupt zip only. MEASURED 2026-09-02 with a chmod 555
            # subfolder: `PermissionError` escaped `scan_report` and a good file
            # sorting after it was never extracted.
            try:
                companion.write_text(content, encoding="utf-8")
            except OSError as exc:
                print(f"  {RED}Failed{RESET}  {companion.name}: "
                      f"{type(exc).__name__}: {exc}")
                failures.append((filepath, f"{type(exc).__name__}: {exc}"))
                continue
            print(f"  {GREEN}Created{RESET}  {companion.name}")
            extracted.append((filepath, companion))
        else:
            print(f"  {RED}Failed{RESET}  Could not extract {filepath.name}")
            failures.append((filepath, "extractor returned nothing"))

    if failures:
        # Named, not swallowed. The batch surviving one bad file is only half
        # the fix; the other half is that the operator can tell a run of 40
        # that produced 40 from a run of 40 that produced 38.
        print(f"\n{YELLOW}{len(failures)} file(s) could not be extracted:{RESET}")
        for path, why in failures:
            print(f"  {YELLOW}{path.name}{RESET}: {why}")

    return ScanResult(extracted, failures)


def update_index(extracted_files):
    """Add newly extracted files to INDEX.md."""
    index = index_file()
    datastore = datastore_dir()
    if not index.exists():
        print(f"{YELLOW}INDEX.md not found - skipping index update{RESET}")
        return

    try:
        content = index.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError) as exc:
        # Same failure the comment below already names, one line earlier: this
        # runs LAST, after every file has been extracted, so raising here
        # reports failure over work that succeeded. `UnicodeDecodeError`
        # subclasses ValueError and this read had no handler at all.
        print(f"{YELLOW}INDEX.md unreadable ({exc}) - skipping index "
              f"update{RESET}")
        return
    today = datetime.now(get_default_tz()).strftime("%Y-%m-%d")

    new_rows = []
    for orig, companion in extracted_files:
        # A target outside `datastore/` is allowed on the CLI, and
        # `relative_to` raised ValueError for it — AFTER every file had already
        # been extracted, so the work was done and the index update died.
        try:
            rel_path = orig.relative_to(datastore)
        except ValueError:
            print(f"{YELLOW}Skipping index row for {orig}: it is outside "
                  f"{datastore}, which INDEX.md rows are relative to.{RESET}")
            continue
        domain = rel_path.parts[0] if rel_path.parts else "unknown"
        # Check if already in index
        if str(rel_path) in content:
            continue
        new_rows.append(f"| `{rel_path}` | {domain.title()} | {orig.stem} (auto-extracted) | {today} | *Review and update validates column* |")

    if not new_rows:
        print(f"{GREEN}INDEX.md already up to date{RESET}")
        return

    # Insert before the LAST HTML comment, not the first one anywhere in the
    # file. The comment above this said "before the closing comment", and the
    # code replaced the first `<!--` it found — so an INDEX.md opening with a
    # header note or a licence comment got data rows injected above it, before
    # the table they belong to.
    marker = content.rfind("<!--")
    block = "\n".join(new_rows)
    if marker != -1:
        content = content[:marker] + block + "\n\n" + content[marker:]
    else:
        content = content.rstrip("\n") + "\n" + block + "\n"

    # Rewrite the whole "Last updated" LINE, and carry its old value into a
    # single Previous line. Replacing only the literal prefix left the old date
    # dangling after `> Previous: `, and the next run did it again — nesting
    # one more stale Previous into the block every time.
    content = re.sub(
        r"^> Last updated:[^\n]*(?:\n> Previous:[^\n]*)?",
        lambda m: (f"> Last updated: {today}  \n> Previous: "
                   f"{m.group(0).split(':', 1)[1].split(chr(10))[0].strip()}"),
        content,
        count=1,
        flags=re.MULTILINE,
    )

    index.write_text(content, encoding="utf-8")
    print(f"{GREEN}Added {len(new_rows)} entries to INDEX.md{RESET}")


def main():
    parser = argparse.ArgumentParser(description="31C DataStore Extraction")
    parser.add_argument("target", nargs="?", default=None,
                        help="Specific directory to scan (default: entire datastore/)")
    parser.add_argument("--update-index", action="store_true",
                        help="Update INDEX.md with new entries")
    parser.add_argument("--force", action="store_true",
                        help="Re-extract even if companion file exists")
    args = parser.parse_args()

    print(f"\n{BOLD}31C DataStore Extraction{RESET}")
    print(f"DataStore: {datastore_dir()}\n")

    result = scan_report(args.target, force=args.force)

    if result.extracted and args.update_index:
        update_index(result.extracted)

    print(f"\n{BOLD}Done.{RESET} Extracted {len(result.extracted)} file(s).")

    # A run that could not read some of its input is not a successful run. The
    # failures were already named above; this is what lets a caller that is not
    # reading the output - a timer, a hook, a shell `&&` - tell the difference.
    if result.failures:
        print(f"{RED}{len(result.failures)} file(s) failed to extract.{RESET}")
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
